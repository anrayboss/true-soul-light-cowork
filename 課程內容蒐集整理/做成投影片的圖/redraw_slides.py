import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Theme Colors (Cyber Neon)
BG_COLOR = RGBColor(11, 15, 25)          # Dark Navy
CARD_BG = RGBColor(21, 29, 45)           # Glass Navy
BORDER_COLOR = RGBColor(129, 140, 248)    # Indigo
TITLE_COLOR = RGBColor(192, 132, 252)     # Violet/Purple
SUBTITLE_COLOR = RGBColor(129, 140, 248)  # Indigo
TEXT_COLOR = RGBColor(241, 245, 249)      # Off-white
HIGHLIGHT_COLOR = RGBColor(244, 114, 182) # Pink
GREEN_COLOR = RGBColor(52, 211, 153)      # Teal
RED_COLOR = RGBColor(239, 68, 68)         # Red

FONT_NAME = 'Microsoft JhengHei'

def apply_slide_base(prs, slide):
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Outer thin border
    border_margin = Inches(0.1)
    border = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        border_margin, border_margin,
        prs.slide_width - 2 * border_margin,
        prs.slide_height - 2 * border_margin
    )
    border.fill.background()
    border.line.color.rgb = BORDER_COLOR
    border.line.width = Pt(0.5)
    if border.adjustments:
        border.adjustments[0] = 0.02

def add_slide_header(slide, title, subtitle=None):
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_NAME
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    
    # Subtitle
    if subtitle:
        p_sub = tf.add_paragraph()
        run_sub = p_sub.add_run()
        run_sub.text = subtitle
        run_sub.font.name = FONT_NAME
        run_sub.font.size = Pt(14)
        run_sub.font.color.rgb = SUBTITLE_COLOR
        
    # Underline
    line_y = Inches(1.3) if not subtitle else Inches(1.5)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), line_y, Inches(11.733), Pt(1.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER_COLOR
    line.line.fill.background()

def add_card(slide, title, bullets, left, top, width, height, accent_color=BORDER_COLOR):
    # Base rounded card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = accent_color
    card.line.width = Pt(1.5)
    if card.adjustments:
        card.adjustments[0] = 0.04
        
    # Text content
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Card Title
    p_title = tf.paragraphs[0]
    p_title.space_after = Pt(10)
    run_title = p_title.add_run()
    run_title.text = title
    run_title.font.name = FONT_NAME
    run_title.font.size = Pt(18)
    run_title.font.bold = True
    run_title.font.color.rgb = accent_color
    
    # Bullets
    for i, b in enumerate(bullets):
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.level = 0
        
        # Check if it starts with a number or standard bullet
        run_bullet = p.add_run()
        if b.startswith("- ") or b.startswith("* "):
            run_bullet.text = "✦ "
            b_text = b[2:]
        else:
            run_bullet.text = ""
            b_text = b
            
        run_bullet.font.name = FONT_NAME
        run_bullet.font.size = Pt(13)
        run_bullet.font.bold = True
        run_bullet.font.color.rgb = accent_color
        
        run_text = p.add_run()
        run_text.text = b_text
        run_text.font.name = FONT_NAME
        run_text.font.size = Pt(13)
        run_text.font.color.rgb = TEXT_COLOR
        
        # Sub-bullets support
        if "\n" in b_text:
            lines = b_text.split("\n")
            p.text = "" # Clear main text to draw manually
            
            run_bullet = p.add_run()
            run_bullet.text = "✦ "
            run_bullet.font.name = FONT_NAME
            run_bullet.font.size = Pt(13)
            run_bullet.font.color.rgb = accent_color
            
            run_text = p.add_run()
            run_text.text = lines[0]
            run_text.font.name = FONT_NAME
            run_text.font.size = Pt(13)
            run_text.font.color.rgb = TEXT_COLOR
            
            for sub_line in lines[1:]:
                p_sub = tf.add_paragraph()
                p_sub.left_indent = Inches(0.3)
                p_sub.space_before = Pt(2)
                p_sub.space_after = Pt(2)
                
                run_sub_bullet = p_sub.add_run()
                run_sub_bullet.text = "• "
                run_sub_bullet.font.name = FONT_NAME
                run_sub_bullet.font.size = Pt(11)
                run_sub_bullet.font.color.rgb = SUBTITLE_COLOR
                
                run_sub_text = p_sub.add_run()
                run_sub_text.text = sub_line.strip()
                run_sub_text.font.name = FONT_NAME
                run_sub_text.font.size = Pt(11)
                run_sub_text.font.color.rgb = TEXT_COLOR

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: 選題庫 x 新手定位賽道
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide1)
    add_slide_header(slide1, "選題庫 x 新手定位賽道", "用熱門主題找到內容方向，再從定位賽道建立長期輸出")
    
    # Two large columns
    left_bullets = [
        "九大熱門主題，抓住風口與人性的切角",
        "01 生活日常 - 生活點滴、趣味日常",
        "02 家居生活 - 家居佈置、收納整理",
        "03 美妝穿搭 - 彩妝教學、穿衣搭配",
        "04 寵物萌寵 - 寵物生活、照顧心得",
        "05 美食旅遊 - 餐廳推薦、旅遊行程",
        "06 娛樂搞笑 - 趣事分享、幽默段子",
        "07 時尚美學 - 時尚趨勢、穿搭品味",
        "08 成長自律 - 自我提升、學習方法",
        "09 理財搞錢 - 理財規劃、副業變現"
    ]
    add_card(slide1, "熱門選題庫", left_bullets, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0), TITLE_COLOR)
    
    right_bullets = [
        "找出你可以持續產出的內容利基",
        "01 關聯定位\n結合多重主題的交集點",
        "02 專業定位\n展現垂直領域的知識權威",
        "03 娛樂定位\n提供趣味與輕鬆的內容",
        "04 吐槽定位\n用真實與批判性引發共鳴",
        "05 乾貨定位\n提供具體可執行的工具/方法",
        "06 美學定位\n重視視覺美感與生活品味",
        "07 故事定位\n以個人經歷與情商建立認同",
        "08 挑戰定位\n設定目標並記錄執行過程",
        "09 避坑定位\n幫助受眾防範常見錯誤",
        "10 跨界定位\n將兩個不相關領域融合",
        "11 資源定位\n做為資訊整理與管道整合者",
        "12 陪伴定位\n建立深度陪伴的社群體驗"
    ]
    # For right column, let's make the card taller or smaller font if needed, let's adjust height.
    # In python-pptx we will let text frame overflow or set smaller font sizes
    add_card(slide1, "新手定位賽道", right_bullets, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.0), GREEN_COLOR)

    # ==========================================
    # SLIDE 2: 市場利基四象限
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide2)
    add_slide_header(slide2, "市場利基四象限", "用 需求 x 競爭 找出黃金縫隙")
    
    # 2x2 grid representing quadrants
    # Top-Left: High Demand, Low Comp -> Golden Niche
    add_card(slide2, "高需求 x 低競爭 (黃金縫隙)", [
        "有需求但還沒被做好 = 最賺錢的利基",
        "用專業差異切入，做高信任 x 高單價",
        "掌握定價能力，脫離價格戰競爭"
    ], Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3), GREEN_COLOR)
    
    # Top-Right: High Demand, High Comp -> Red Ocean
    add_card(slide2, "高需求 x 高競爭 (紅海市場)", [
        "市場很大，但一進去就容易被淹沒",
        "解法：深度細分 or 重新定位「切一刀」",
        "例：減重改成「40歲婦女減脂」",
        "例：皮膚美容改成「痘痘肌專門處理」"
    ], Inches(6.9), Inches(1.8), Inches(5.6), Inches(2.3), RED_COLOR)
    
    # Bottom-Left: Low Demand, Low Comp -> Observation Area
    add_card(slide2, "低需求 x 低競爭 (觀察測試區)", [
        "現在沒市場，不代表未來沒有",
        "主要任務：教育市場 + 測試內容方向",
        "例：AI個人分身教練、情緒價值陪伴IP"
    ], Inches(0.8), Inches(4.4), Inches(5.6), Inches(2.3), BORDER_COLOR)
    
    # Bottom-Right: Low Demand, High Comp -> Black Hole
    add_card(slide2, "低需求 x 高競爭 (流量黑洞)", [
        "避開！照抄型分享且無定位，什麼都發",
        "沒差異 x 沒需求 = 做再多內容也無法變現",
        "沒有差異的技能，最後只能打價格戰"
    ], Inches(6.9), Inches(4.4), Inches(5.6), Inches(2.3), RGBColor(100, 116, 139))

    # Axis label at the bottom center
    axis_box = slide2.shapes.add_textbox(Inches(5.0), Inches(6.8), Inches(3.333), Inches(0.4))
    axis_tf = axis_box.text_frame
    axis_p = axis_tf.paragraphs[0]
    axis_p.alignment = PP_ALIGN.CENTER
    axis_run = axis_p.add_run()
    axis_run.text = "← 競爭者密度 →"
    axis_run.font.name = FONT_NAME
    axis_run.font.size = Pt(14)
    axis_run.font.bold = True
    axis_run.font.color.rgb = SUBTITLE_COLOR

    # ==========================================
    # SLIDE 3: 定位調整的五個步驟
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide3)
    add_slide_header(slide3, "別怪市場不好，是你站錯位置", "定位調整的五個具體步驟")
    
    steps = [
        ("STEP 1", "盤點現況", "盤點自己目前在哪一個四象限格子中"),
        ("STEP 2", "紅海細分", "如果在紅海，立即做深度細分切一刀"),
        ("STEP 3", "黑洞重組", "如果在黑洞，必須重做定位與技能重組"),
        ("STEP 4", "觀察測試", "如果在觀察區，持續測試內容並教育市場"),
        ("STEP 5", "移動目標", "終極目標：移動到「黃金縫隙」賺取高利潤")
    ]
    
    for idx, (step_num, title, desc) in enumerate(steps):
        y_pos = Inches(1.8 + idx * 1.0)
        
        # Step label shape
        label = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(2.0), Inches(0.8))
        label.fill.solid()
        label.fill.fore_color.rgb = HIGHLIGHT_COLOR if idx == 4 else CARD_BG
        label.line.color.rgb = TITLE_COLOR
        label.line.width = Pt(1)
        
        label_tf = label.text_frame
        label_p = label_tf.paragraphs[0]
        label_p.alignment = PP_ALIGN.CENTER
        label_run = label_p.add_run()
        label_run.text = step_num
        label_run.font.name = FONT_NAME
        label_run.font.size = Pt(16)
        label_run.font.bold = True
        label_run.font.color.rgb = TEXT_COLOR
        
        # Description box
        desc_box = slide3.shapes.add_textbox(Inches(3.0), y_pos, Inches(9.5), Inches(0.8))
        desc_tf = desc_box.text_frame
        desc_tf.word_wrap = True
        
        p = desc_tf.paragraphs[0]
        run_title = p.add_run()
        run_title.text = title + "  —  "
        run_title.font.name = FONT_NAME
        run_title.font.size = Pt(18)
        run_title.font.bold = True
        run_title.font.color.rgb = TITLE_COLOR
        
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.name = FONT_NAME
        run_desc.font.size = Pt(15)
        run_desc.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 4: 信任內容的證據矩陣 (火力展示)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide4)
    add_slide_header(slide4, "信任內容的證據矩陣", "火力展示：多維度建立受眾信任")
    
    cols = [
        ("客戶成果", "目的：顯示結果", ["- 前後對照圖 (對比)", "- 學員故事 (親身歷程)"], GREEN_COLOR),
        ("社會認可", "目的：顯示公信力", ["- 媒體報導 (第三方)", "- 受邀講座 (權威度)"], TITLE_COLOR),
        ("口碑傳遞", "目的：顯示影響力", ["- 推薦文 (使用者推薦)", "- UGC 真實分享 (社群擴散)"], BORDER_COLOR)
    ]
    
    for idx, (title, purpose, items, color) in enumerate(cols):
        left_pos = Inches(0.8 + idx * 3.9)
        # Main Column Card
        add_card(slide4, title, items, left_pos, Inches(1.8), Inches(3.7), Inches(3.5), color)
        
        # Bottom Purpose Box
        p_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(5.6), Inches(3.7), Inches(1.0))
        p_box.fill.solid()
        p_box.fill.fore_color.rgb = CARD_BG
        p_box.line.color.rgb = color
        p_box.line.width = Pt(1)
        p_tf = p_box.text_frame
        p_p = p_tf.paragraphs[0]
        p_p.alignment = PP_ALIGN.CENTER
        p_run = p_p.add_run()
        p_run.text = purpose
        p_run.font.name = FONT_NAME
        p_run.font.size = Pt(16)
        p_run.font.bold = True
        p_run.font.color.rgb = color

    # ==========================================
    # SLIDE 5: 公域 vs. 私域
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide5)
    add_slide_header(slide5, "公域 vs. 私域", "公域引流獲取注意，私域精準留存變現")
    
    add_card(slide5, "公域流量 (大海)", [
        "管道：IG、YouTube、Facebook、其他自媒體等",
        "特點：受演算法控制，流量不屬於你",
        "風險：平台規則一改，觸及率可能瞬間腰斬",
        "核心價值：做為引流與漏斗最上層"
    ], Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), TITLE_COLOR)
    
    add_card(slide5, "私域留量 (魚池)", [
        "管道：Email 電子報、LINE 官方帳號、電話名單",
        "特點：主動觸及，不需平台允許，沒有演算法限制",
        "優勢：不懼演算法改版，與用戶建立深度高信任連結",
        "核心價值：做為高信任維護與核心成交變現場景"
    ], Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8), GREEN_COLOR)

    # ==========================================
    # SLIDE 6: 這堂陪跑課，就是想改變你
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide6)
    add_slide_header(slide6, "這堂陪跑課，就是想改變你", "突破瓶頸，打造個人品牌的質變工程")
    
    changes = [
        ("1", "定位不模糊", "找到熱愛/擅長且能賺錢的方向，確立此生靈魂賽道"),
        ("2", "脫離窮忙困境", "設計高信任成交模式，不再依賴低單價勞力變現"),
        ("3", "自媒體全方位獲利", "自媒體有流量、更有「留量」與影響力，並有變現成績，成為指標"),
        ("4", "突破時間困境", "透過 AI 自動化工具降本增效，實現「睡覺也能成交」的自動化系統"),
        ("5", "創造斜槓收入", "創造持續性收入，賺更多錢的同時，也能用你的專業幫助更多人"),
        ("6", "獲得真正的自由", "重獲注意力、時間、金錢、健康與幸福的掌控權")
    ]
    
    for idx, (num, title, desc) in enumerate(changes):
        y_pos = Inches(1.8 + idx * 0.85)
        
        # Circle badge
        badge = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y_pos, Inches(0.6), Inches(0.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = HIGHLIGHT_COLOR
        badge.line.fill.background()
        
        b_tf = badge.text_frame
        b_p = b_tf.paragraphs[0]
        b_p.alignment = PP_ALIGN.CENTER
        b_run = b_p.add_run()
        b_run.text = num
        b_run.font.name = FONT_NAME
        b_run.font.size = Pt(14)
        b_run.font.bold = True
        b_run.font.color.rgb = TEXT_COLOR
        
        # Text Callout
        tb = slide6.shapes.add_textbox(Inches(1.6), y_pos - Inches(0.05), Inches(10.9), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        
        run_title = p.add_run()
        run_title.text = title + "  —  "
        run_title.font.name = FONT_NAME
        run_title.font.size = Pt(16)
        run_title.font.bold = True
        run_title.font.color.rgb = TITLE_COLOR
        
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.name = FONT_NAME
        run_desc.font.size = Pt(14)
        run_desc.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 7: 納瓦爾：注意力 > 時間 > 金錢
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide7)
    add_slide_header(slide7, "納瓦爾：注意力 > 時間 > 金錢", "真正稀缺的不是資源，而是你持續吸引注意力的能力")
    
    navar_cols = [
        ("01  注意力", ["最稀缺的槓桿資源", "先被看見，才有機會被選擇", "是個人品牌的起點"], TITLE_COLOR),
        ("02  時間", ["每個人一天都只有24小時", "關鍵在於把時間放在高價值事務上", "透過系統與自動化放大時間產值"], BORDER_COLOR),
        ("03  金錢", ["通常是價值的結果而非起點", "能持續創造並傳遞價值的人", "最後會吸走市場上最多的金錢資源"], GREEN_COLOR)
    ]
    
    for idx, (title, items, color) in enumerate(navar_cols):
        left_pos = Inches(0.8 + idx * 3.9)
        add_card(slide7, title, items, left_pos, Inches(1.8), Inches(3.7), Inches(3.4), color)
        
    # Bottom Summary
    summary_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.733), Inches(1.2))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = CARD_BG
    summary_box.line.color.rgb = HIGHLIGHT_COLOR
    summary_box.line.width = Pt(2)
    s_tf = summary_box.text_frame
    s_p = s_tf.paragraphs[0]
    s_p.alignment = PP_ALIGN.CENTER
    s_run = s_p.add_run()
    s_run.text = "✨ 先搶佔注意力，再放大時間價值，金錢才會順理成章跟上。 ✨"
    s_run.font.name = FONT_NAME
    s_run.font.size = Pt(20)
    s_run.font.bold = True
    s_run.font.color.rgb = HIGHLIGHT_COLOR

    # ==========================================
    # SLIDE 8: 誰要來上這堂課程？
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide8)
    add_slide_header(slide8, "誰要來上這堂課程？", "適合不同發展階段的個人與組織")
    
    targets = [
        ("老闆 / 創業家", ["知道方向但一定要做到", "思維提升，團隊才能跟上", "建立企業長線的流量與變現系統"], TITLE_COLOR),
        ("品牌主管 / 行銷小編", ["協助公司品牌更上一層樓", "掌握自媒體最新的實戰策略", "以成效爭取升職加薪"], BORDER_COLOR),
        ("專業人士 / 療癒者", ["醫師、講師、身心靈專家等", "將專業無形價值產品化、數據化", "擺脫一對一勞力變現上限"], GREEN_COLOR),
        ("內容操盤手 / 新手IP", ["未來想透過教學/打造他人變現", "希望建立副業與斜槓收入", "此課程將是你的實戰起點"], HIGHLIGHT_COLOR)
    ]
    
    for idx, (title, items, color) in enumerate(targets):
        col = idx % 2
        row = idx // 2
        left_pos = Inches(0.8 + col * 5.9)
        top_pos = Inches(1.8 + row * 2.5)
        add_card(slide8, title, items, left_pos, top_pos, Inches(5.6), Inches(2.2), color)

    # ==========================================
    # SLIDE 9: 精銳 IP 的 5 大資產
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide9)
    add_slide_header(slide9, "精銳 IP 的 5 大資產", "個人品牌的底層資產積累，決定變現的可持續性")
    
    assets = [
        ("1. 定位資產", "釐清天賦與熱愛，切入利基黃金縫隙"),
        ("2. 信任資產", "累積學員案例、公信力與口碑傳遞"),
        ("3. 留量資產", "建立不畏演算法的私域流量魚池"),
        ("4. 系統資產", "打造自動化成交與降本增效流程"),
        ("5. 延伸資產", "擴展代理商、產品線與跨界合作機會")
    ]
    
    for idx, (title, desc) in enumerate(assets):
        left_pos = Inches(0.8 + idx * 2.38)
        
        # Vertical card
        card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.2), Inches(2.2), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = TITLE_COLOR
        card.line.width = Pt(1.5)
        
        c_tf = card.text_frame
        c_tf.word_wrap = True
        
        p_title = c_tf.paragraphs[0]
        p_title.alignment = PP_ALIGN.CENTER
        p_title.space_after = Pt(14)
        run_title = p_title.add_run()
        run_title.text = title
        run_title.font.name = FONT_NAME
        run_title.font.size = Pt(18)
        run_title.font.bold = True
        run_title.font.color.rgb = TITLE_COLOR
        
        p_desc = c_tf.add_paragraph()
        p_desc.alignment = PP_ALIGN.CENTER
        run_desc = p_desc.add_run()
        run_desc.text = desc
        run_desc.font.name = FONT_NAME
        run_desc.font.size = Pt(13)
        run_desc.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 10: 你現在是哪一種狀態？
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide10)
    add_slide_header(slide10, "你現在是哪一種狀態？", "信任系統升級路徑：從被動曝光到口碑裂變")
    
    states = [
        ("路過型", ["有曝光但沒成交", "缺乏信任建立與明確呼籲"], BORDER_COLOR),
        ("猶豫型", ["有人詢問但遲不下單", "產品感知價值不夠，或疑慮未消"], RGBColor(245, 158, 11)), # Amber
        ("硬推型", ["會成交但過程非常累", "缺乏自動化跟進系統，全靠銷售技巧"], RED_COLOR),
        ("信任型", ["有回購，且客戶主動介紹", "最理想狀態，累積高黏性私域資產"], GREEN_COLOR)
    ]
    
    for idx, (title, items, color) in enumerate(states):
        col = idx % 2
        row = idx // 2
        left_pos = Inches(0.8 + col * 5.9)
        top_pos = Inches(1.8 + row * 2.5)
        add_card(slide10, title, items, left_pos, top_pos, Inches(5.6), Inches(2.2), color)

    # ==========================================
    # SLIDE 11: 客戶為什麼不選你？
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide11)
    add_slide_header(slide11, "客戶為什麼不選你？", "直擊成交的核心痛點")
    
    quote_card = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(10.333), Inches(4.0))
    quote_card.fill.solid()
    quote_card.fill.fore_color.rgb = CARD_BG
    quote_card.line.color.rgb = RED_COLOR
    quote_card.line.width = Pt(3)
    
    q_tf = quote_card.text_frame
    q_tf.word_wrap = True
    q_tf.margin_top = Inches(0.8)
    
    p = q_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)
    run1 = p.add_run()
    run1.text = "「 因為你沒有系統和方法！ 」"
    run1.font.name = FONT_NAME
    run1.font.size = Pt(36)
    run1.font.bold = True
    run1.font.color.rgb = HIGHLIGHT_COLOR
    
    p2 = q_tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "所有厲害的人，都在透過「系統」和「方法」來加速前進。\n擺脫個人體力限制，才能實現長久而穩定的個人品牌變現。"
    run2.font.name = FONT_NAME
    run2.font.size = Pt(20)
    run2.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 12: 每個階段的關鍵流程
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide12)
    add_slide_header(slide12, "每個階段的關鍵流程", "明確劃分階段任務，直達最終變現")
    
    flow_steps = [
        ("初期", "建立人設", "精準定位，傳遞個人風格"),
        ("中期", "篩選人群", "精準吸客，排除非受眾"),
        ("後期", "累積信任", "交付價值，打消購買疑慮")
    ]
    
    for idx, (stage, title, desc) in enumerate(flow_steps):
        left_pos = Inches(0.8 + idx * 3.0)
        
        # Step card
        card = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(2.2), Inches(2.6), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = TITLE_COLOR
        card.line.width = Pt(1.5)
        
        c_tf = card.text_frame
        c_tf.word_wrap = True
        
        p1 = c_tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.space_after = Pt(10)
        r1 = p1.add_run()
        r1.text = stage
        r1.font.name = FONT_NAME
        r1.font.size = Pt(16)
        r1.font.color.rgb = SUBTITLE_COLOR
        
        p2 = c_tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_after = Pt(10)
        r2 = p2.add_run()
        r2.text = title
        r2.font.name = FONT_NAME
        r2.font.size = Pt(22)
        r2.font.bold = True
        r2.font.color.rgb = TITLE_COLOR
        
        p3 = c_tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = desc
        r3.font.name = FONT_NAME
        r3.font.size = Pt(13)
        r3.font.color.rgb = TEXT_COLOR
        
        # Draw Arrow indicator to next step (if not last)
        if idx < 2:
            arrow = slide12.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left_pos + Inches(2.65), Inches(3.6), Inches(0.3), Inches(0.5))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BORDER_COLOR
            arrow.line.fill.background()

    # Destination Box (變現)
    dest_box = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(2.2), Inches(2.7), Inches(3.5))
    dest_box.fill.solid()
    dest_box.fill.fore_color.rgb = CARD_BG
    dest_box.line.color.rgb = GREEN_COLOR
    dest_box.line.width = Pt(2.5)
    
    d_tf = dest_box.text_frame
    d_tf.word_wrap = True
    d_tf.margin_top = Inches(0.8)
    
    dp1 = d_tf.paragraphs[0]
    dp1.alignment = PP_ALIGN.CENTER
    dp1.space_after = Pt(10)
    dr1 = dp1.add_run()
    dr1.text = "最終目標"
    dr1.font.name = FONT_NAME
    dr1.font.size = Pt(16)
    dr1.font.color.rgb = GREEN_COLOR
    
    dp2 = d_tf.add_paragraph()
    dp2.alignment = PP_ALIGN.CENTER
    dr2 = dp2.add_run()
    dr2.text = "變現\n(好服務與好產品)"
    dr2.font.name = FONT_NAME
    dr2.font.size = Pt(20)
    dr2.font.bold = True
    dr2.font.color.rgb = GREEN_COLOR

    # ==========================================
    # SLIDE 13: 急著賺錢的人，很容易賺不到錢
    # ==========================================
    slide13 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide13)
    add_slide_header(slide13, "急著賺錢的人，很容易賺不到錢", "先長根，再長大：精銳 IP 的反向邏輯")
    
    add_card(slide13, "雜草 (短線思維)", [
        "急著賺錢變現，缺乏耐心",
        "沒有清晰定位，內容發散",
        "信任還沒建立就急於銷售",
        "一下想變現，一下又換方向"
    ], Inches(0.8), Inches(1.8), Inches(5.6), Inches(3.2), RED_COLOR)
    
    add_card(slide13, "大樹 (長線思維)", [
        "先站穩位置，打好定位根基",
        "持續輸出價值，建立深厚信任",
        "慢慢放大自媒體的長線影響力",
        "根基穩固，最後自然穩定變現"
    ], Inches(6.9), Inches(1.8), Inches(5.6), Inches(3.2), GREEN_COLOR)
    
    # Bottom Flow
    flow_box = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.3), Inches(11.733), Inches(1.4))
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = CARD_BG
    flow_box.line.color.rgb = HIGHLIGHT_COLOR
    flow_box.line.width = Pt(1.5)
    f_tf = flow_box.text_frame
    
    fp1 = f_tf.paragraphs[0]
    fp1.alignment = PP_ALIGN.CENTER
    fp1.space_after = Pt(6)
    fr1 = fp1.add_run()
    fr1.text = "精銳 IP 反向變現邏輯"
    fr1.font.name = FONT_NAME
    fr1.font.size = Pt(16)
    fr1.font.bold = True
    fr1.font.color.rgb = HIGHLIGHT_COLOR
    
    fp2 = f_tf.add_paragraph()
    fp2.alignment = PP_ALIGN.CENTER
    fr2 = fp2.add_run()
    fr2.text = "1) 站穩位置  →  2) 建立信任  →  3) 放大影響力  →  4) 穩定變現"
    fr2.font.name = FONT_NAME
    fr2.font.size = Pt(18)
    fr2.font.bold = True
    fr2.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 14: 第三條錯路：信任還沒開始，你就急著要成交
    # ==========================================
    slide14 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide14)
    add_slide_header(slide14, "第三條錯路：信任還沒開始，你就急著要成交", "忽視認知階段，強推只會適得其反")
    
    # Left column: Progression
    prog_box = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    prog_box.fill.solid()
    prog_box.fill.fore_color.rgb = CARD_BG
    prog_box.line.color.rgb = RED_COLOR
    prog_box.line.width = Pt(2)
    p_tf = prog_box.text_frame
    p_tf.word_wrap = True
    p_tf.margin_top = Inches(0.4)
    
    p_items = [
        "客戶眼中的錯誤推銷流程：",
        "❌ 尚未被看見  (曝光不足)",
        "❌ 尚未被理解  (價值模糊)",
        "❌ 尚未被記住  (印象淺薄)",
        "🚨 馬上就急著變現 (BUY NOW)"
    ]
    for idx, item in enumerate(p_items):
        p = p_tf.paragraphs[0] if idx == 0 else p_tf.add_paragraph()
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = item
        run.font.name = FONT_NAME
        run.font.size = Pt(18)
        if idx == 0:
            run.font.bold = True
            run.font.color.rgb = RED_COLOR
        elif idx == 4:
            run.font.bold = True
            run.font.color.rgb = HIGHLIGHT_COLOR
        else:
            run.font.color.rgb = TEXT_COLOR

    # Right column: Metaphor card
    metaphor_box = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    metaphor_box.fill.solid()
    metaphor_box.fill.fore_color.rgb = CARD_BG
    metaphor_box.line.color.rgb = TITLE_COLOR
    metaphor_box.line.width = Pt(2)
    m_tf = metaphor_box.text_frame
    m_tf.word_wrap = True
    m_tf.margin_top = Inches(1.0)
    
    mp1 = m_tf.paragraphs[0]
    mp1.alignment = PP_ALIGN.CENTER
    mp1.space_after = Pt(20)
    mr1 = mp1.add_run()
    mr1.text = "💡 關係比喻 💡"
    mr1.font.name = FONT_NAME
    mr1.font.size = Pt(22)
    mr1.font.bold = True
    mr1.font.color.rgb = TITLE_COLOR
    
    mp2 = m_tf.add_paragraph()
    mp2.alignment = PP_ALIGN.CENTER
    mp2.space_after = Pt(16)
    mr2 = mp2.add_run()
    mr2.text = "「 牽手都還沒，你就先想結婚。 」"
    mr2.font.name = FONT_NAME
    mr2.font.size = Pt(20)
    mr2.font.bold = True
    mr2.font.color.rgb = HIGHLIGHT_COLOR
    
    mp3 = m_tf.add_paragraph()
    mp3.alignment = PP_ALIGN.CENTER
    mr3 = mp3.add_run()
    mr3.text = "節奏太快、缺乏信任累積，\n對方只會本能地後退與拒絕。"
    mr3.font.name = FONT_NAME
    mr3.font.size = Pt(16)
    mr3.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 15: 客戶不是想變專家，他只是想不要選錯
    # ==========================================
    slide15 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide15)
    add_slide_header(slide15, "客戶不是想變專家，他只是想不要選錯", "換位思考：從專家視角轉向客戶視角")
    
    add_card(slide15, "專家視角 (溝通阻礙)", [
        "❌ 一味灌輸學術知識",
        "❌ 滿口講深奧大理論",
        "❌ 說教感太重，顯得高高在上"
    ], Inches(0.8), Inches(1.8), Inches(5.6), Inches(3.0), RED_COLOR)
    
    add_card(slide15, "客戶視角 (真實需求)", [
        "🛡️ 用戶本能怕出錯",
        "🛡️ 用戶想規避選擇風險",
        "🛡️ 用戶恐懼失敗的代價"
    ], Inches(6.9), Inches(1.8), Inches(5.6), Inches(3.0), GREEN_COLOR)
    
    # Highlight & Translation callout
    bottom_card = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.1), Inches(11.733), Inches(1.6))
    bottom_card.fill.solid()
    bottom_card.fill.fore_color.rgb = CARD_BG
    bottom_card.line.color.rgb = HIGHLIGHT_COLOR
    bottom_card.line.width = Pt(2)
    b_tf = bottom_card.text_frame
    
    bp1 = b_tf.paragraphs[0]
    bp1.alignment = PP_ALIGN.CENTER
    bp1.space_after = Pt(8)
    br1 = bp1.add_run()
    br1.text = "看到結果，才是他們真正想要的！"
    br1.font.name = FONT_NAME
    br1.font.size = Pt(20)
    br1.font.bold = True
    br1.font.color.rgb = HIGHLIGHT_COLOR
    
    bp2 = b_tf.add_paragraph()
    bp2.alignment = PP_ALIGN.CENTER
    br2 = bp2.add_run()
    br2.text = "關鍵在於： 將你的「專業語言」 ──翻譯為── 顧客聽得懂的「客戶語言」"
    br2.font.name = FONT_NAME
    br2.font.size = Pt(16)
    br2.font.bold = True
    br2.font.color.rgb = TEXT_COLOR

    # Save presentation
    output_path = r"d:\Git\true-soul-light\不用版控的\做成投影片的圖\做成投影片的圖_重繪版.pptx"
    prs.save(output_path)
    print(f"Redrawn presentation saved successfully to: {output_path}")

if __name__ == '__main__':
    main()
