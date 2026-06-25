import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Theme Colors (Cyber Neon)
BG_COLOR = RGBColor(11, 15, 25)          # Dark Navy
CARD_BG = RGBColor(21, 29, 45)           # Glass Navy
BORDER_COLOR = RGBColor(129, 140, 248)    # Indigo
TITLE_COLOR = RGBColor(192, 132, 252)     # Violet/Purple
SUBTITLE_COLOR = RGBColor(129, 140, 248)  # Indigo
TEXT_COLOR = RGBColor(241, 245, 249)      # Off-white
HIGHLIGHT_COLOR = RGBColor(244, 114, 182) # Pink
GREEN_COLOR = RGBColor(52, 211, 153)      # Teal
RED_COLOR = RGBColor(239, 68, 68)         # Red

FONT_NAME = 'Microsoft JhengHei'

def apply_slide_base(prs, slide):
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Outer thin border
    border_margin = Inches(0.1)
    border = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        border_margin, border_margin,
        prs.slide_width - 2 * border_margin,
        prs.slide_height - 2 * border_margin
    )
    border.fill.background()
    border.line.color.rgb = BORDER_COLOR
    border.line.width = Pt(0.5)
    if border.adjustments:
        border.adjustments[0] = 0.02

def add_slide_header(slide, title, subtitle=None):
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_NAME
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    
    # Subtitle
    if subtitle:
        p_sub = tf.add_paragraph()
        run_sub = p_sub.add_run()
        run_sub.text = subtitle
        run_sub.font.name = FONT_NAME
        run_sub.font.size = Pt(14)
        run_sub.font.color.rgb = SUBTITLE_COLOR
        
    # Underline
    line_y = Inches(1.3) if not subtitle else Inches(1.5)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), line_y, Inches(11.733), Pt(1.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER_COLOR
    line.line.fill.background()

def add_card(slide, title, bullets, left, top, width, height, accent_color=BORDER_COLOR):
    # Base rounded card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = accent_color
    card.line.width = Pt(1.5)
    if card.adjustments:
        card.adjustments[0] = 0.04
        
    # Text content
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Card Title
    p_title = tf.paragraphs[0]
    p_title.space_after = Pt(10)
    run_title = p_title.add_run()
    run_title.text = title
    run_title.font.name = FONT_NAME
    run_title.font.size = Pt(18)
    run_title.font.bold = True
    run_title.font.color.rgb = accent_color
    
    # Bullets
    for i, b in enumerate(bullets):
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.level = 0
        
        run_bullet = p.add_run()
        if b.startswith("- ") or b.startswith("* ") or b.startswith("❌ ") or b.startswith("🛡️ "):
            run_bullet.text = b[:2]
            b_text = b[2:]
        else:
            run_bullet.text = "✦ "
            b_text = b
            
        run_bullet.font.name = FONT_NAME
        run_bullet.font.size = Pt(13)
        run_bullet.font.bold = True
        run_bullet.font.color.rgb = accent_color
        
        run_text = p.add_run()
        run_text.text = b_text
        run_text.font.name = FONT_NAME
        run_text.font.size = Pt(13)
        run_text.font.color.rgb = TEXT_COLOR
        
        # Sub-bullets support
        if "\n" in b_text:
            lines = b_text.split("\n")
            p.text = "" # Clear main text to draw manually
            
            run_bullet = p.add_run()
            run_bullet.text = "✦ "
            run_bullet.font.name = FONT_NAME
            run_bullet.font.size = Pt(13)
            run_bullet.font.color.rgb = accent_color
            
            run_text = p.add_run()
            run_text.text = lines[0]
            run_text.font.name = FONT_NAME
            run_text.font.size = Pt(13)
            run_text.font.color.rgb = TEXT_COLOR
            
            for sub_line in lines[1:]:
                p_sub = tf.add_paragraph()
                p_sub.left_indent = Inches(0.3)
                p_sub.space_before = Pt(2)
                p_sub.space_after = Pt(2)
                
                run_sub_bullet = p_sub.add_run()
                run_sub_bullet.text = "• "
                run_sub_bullet.font.name = FONT_NAME
                run_sub_bullet.font.size = Pt(11)
                run_sub_bullet.font.color.rgb = SUBTITLE_COLOR
                
                run_sub_text = p_sub.add_run()
                run_sub_text.text = sub_line.strip()
                run_sub_text.font.name = FONT_NAME
                run_sub_text.font.size = Pt(11)
                run_sub_text.font.color.rgb = TEXT_COLOR

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: 人生三悟
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide1)
    add_slide_header(slide1, "人生三悟", "六維破界系列課程  |  陳老師行銷內容")
    
    add_card(slide1, "悟道 (悟規律)", [
        "了解規律：看透世事運行的底層邏輯",
        "順應自然：不違背自然規律與發展規律",
        "知未來：預見趨勢，提前做局與佈局",
        "知因果：明晰每一種選擇背後的代價"
    ], Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.8), TITLE_COLOR)
    
    add_card(slide1, "悟人 (悟人性)", [
        "了解人性：洞察人們內心的動機與慾望",
        "順應人性：以利他思維滿足他人的需求",
        "順應人性：不做反人性的無效社交"
    ], Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.8), GREEN_COLOR)
    
    add_card(slide1, "悟己 (悟天賦)", [
        "了解自己：清晰自己的天賦與能力邊界",
        "明確目標：不盲從，始終向結果前進",
        "克服內耗：與自己和解，不糾結"
    ], Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), HIGHLIGHT_COLOR)

    # ==========================================
    # SLIDE 2: 思維七級
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide2)
    add_slide_header(slide2, "思維七級", "關係的層次決定你的思維維度  |  六維破界")
    
    levels = [
        ("Level 7", "超我思維", "不受物質和社會約束，遵循道德和快樂原則（擁有極強的洞察力、判斷力、智慧與思考）", HIGHLIGHT_COLOR),
        ("Level 6", "社會思維", "相互交往、相互作用、相互聯繫、相互影響的多維度整合能力", TITLE_COLOR),
        ("Level 5", "辯證思維", "根據事物好與壞的兩面性做出客觀的判斷和選擇", BORDER_COLOR),
        ("Level 4", "利益思維", "以利益動機洞察他人，明白「你的價值就是別人利益的動機」", GREEN_COLOR),
        ("Level 3", "立場思維", "能夠站在對方的角度思考，並得出第三方的客觀觀點", TEXT_COLOR),
        ("Level 2", "情緒思維", "感知身體對思維的情緒反應，易受感性干擾", TEXT_COLOR),
        ("Level 1", "感知思維", "未曾經析，純粹憑藉直覺與感知做出即時判斷", TEXT_COLOR)
    ]
    
    for idx, (lvl, name, desc, color) in enumerate(levels):
        y_pos = Inches(1.8 + idx * 0.72)
        
        # Level Badge
        badge = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(1.8), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = CARD_BG
        badge.line.color.rgb = color
        badge.line.width = Pt(1.5)
        
        b_tf = badge.text_frame
        b_p = b_tf.paragraphs[0]
        b_p.alignment = PP_ALIGN.CENTER
        b_run = b_p.add_run()
        b_run.text = lvl
        b_run.font.name = FONT_NAME
        b_run.font.size = Pt(14)
        b_run.font.bold = True
        b_run.font.color.rgb = color
        
        # Label & Desc
        tb = slide2.shapes.add_textbox(Inches(2.8), y_pos - Inches(0.05), Inches(9.7), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        
        run_title = p.add_run()
        run_title.text = name + "  —  "
        run_title.font.name = FONT_NAME
        run_title.font.size = Pt(15)
        run_title.font.bold = True
        run_title.font.color.rgb = color
        
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.name = FONT_NAME
        run_desc.font.size = Pt(13)
        run_desc.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 3: 思考七級
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide3)
    add_slide_header(slide3, "思考七級", "你思考的深度，決定你的人生高度  |  六維破界")
    
    think_levels = [
        ("6 級", "關鍵因素", "找到關鍵節點，從根本上解決問題（解決根本）", HIGHLIGHT_COLOR),
        ("5 級", "複雜歸因", "多重分析方法結合，全面權衡各種因素", TITLE_COLOR),
        ("4 級", "簡單歸因", "分析原因，尋找對應的解決方法", BORDER_COLOR),
        ("3 級", "結論判斷", "透過結果思考底層邏輯，建立個人結論", GREEN_COLOR),
        ("2 級", "數據與事實", "通過真實數據與客觀事實思考問題", TEXT_COLOR),
        ("1 級", "關注表象", "只看表面，容易被表面現象蒙蔽與干擾", TEXT_COLOR),
        ("0 級", "茫然無知", "缺乏思考框架與感知，對問題茫然無知", TEXT_COLOR)
    ]
    
    for idx, (lvl, name, desc, color) in enumerate(think_levels):
        y_pos = Inches(1.8 + idx * 0.72)
        
        # Level Badge
        badge = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(1.8), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = CARD_BG
        badge.line.color.rgb = color
        badge.line.width = Pt(1.5)
        
        b_tf = badge.text_frame
        b_p = b_tf.paragraphs[0]
        b_p.alignment = PP_ALIGN.CENTER
        b_run = b_p.add_run()
        b_run.text = lvl
        b_run.font.name = FONT_NAME
        b_run.font.size = Pt(14)
        b_run.font.bold = True
        b_run.font.color.rgb = color
        
        # Label & Desc
        tb = slide3.shapes.add_textbox(Inches(2.8), y_pos - Inches(0.05), Inches(9.7), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        
        run_title = p.add_run()
        run_title.text = name + "  —  "
        run_title.font.name = FONT_NAME
        run_title.font.size = Pt(15)
        run_title.font.bold = True
        run_title.font.color.rgb = color
        
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.name = FONT_NAME
        run_desc.font.size = Pt(13)
        run_desc.font.color.rgb = TEXT_COLOR

    # ==========================================
    # SLIDE 4: 商業思維
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide4)
    add_slide_header(slide4, "商業思維", "從交換體力到交換資源，看透商業變現底層")
    
    add_card(slide4, "賣別人的東西 (普通人 / 營銷)", [
        "啟動成本低：不需背負庫存與研發壓力",
        "反饋來得快：驗證市場需求的速度快",
        "成功率較高：藉助已有產品力進行推廣"
    ], Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.6), GREEN_COLOR)
    
    add_card(slide4, "賣自己的東西 (企業老闆 / 創新)", [
        "啟動成本高：需研發投入與品牌創新",
        "需長期積累：回報週期長，但護城河高",
        "高利潤空間：掌控產品定價權與分配權"
    ], Inches(6.9), Inches(1.8), Inches(5.6), Inches(2.6), TITLE_COLOR)
    
    # Bottom exchange layers
    ex_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.6), Inches(11.733), Inches(2.1))
    ex_box.fill.solid()
    ex_box.fill.fore_color.rgb = CARD_BG
    ex_box.line.color.rgb = HIGHLIGHT_COLOR
    ex_box.line.width = Pt(1.5)
    ex_tf = ex_box.text_frame
    ex_tf.word_wrap = True
    
    ep1 = ex_tf.paragraphs[0]
    ep1.space_before = Pt(6)
    ep1.space_after = Pt(10)
    ep1.alignment = PP_ALIGN.CENTER
    er1 = ep1.add_run()
    er1.text = "💡 價 值 交 換 的 三 個 層 次 💡"
    er1.font.name = FONT_NAME
    er1.font.size = Pt(16)
    er1.font.bold = True
    er1.font.color.rgb = HIGHLIGHT_COLOR
    
    ep2 = ex_tf.add_paragraph()
    ep2.alignment = PP_ALIGN.CENTER
    ep2.space_after = Pt(8)
    er2 = ep2.add_run()
    er2.text = "上層：交換權力與資源  |  中層：交換腦力與技能  |  底層：交換體力與時間"
    er2.font.name = FONT_NAME
    er2.font.size = Pt(15)
    er2.font.bold = True
    er2.font.color.rgb = TEXT_COLOR
    
    ep3 = ex_tf.add_paragraph()
    ep3.alignment = PP_ALIGN.CENTER
    er3 = ep3.add_run()
    er3.text = "💎 商業獲利公式： 創造價值  →  告知價值  →  提供價值  =  賺錢 💎"
    er3.font.name = FONT_NAME
    er3.font.size = Pt(14)
    er3.font.bold = True
    er3.font.color.rgb = GREEN_COLOR

    # ==========================================
    # SLIDE 5: 布局高手
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide5)
    add_slide_header(slide5, "布局高手", "高手必備的素質、能力與思考架構")
    
    add_card(slide5, "6 大素質", [
        "熱愛工作 - 熱情是持續的關鍵",
        "積極樂觀 - 面對挫折的心態",
        "明確目標 - 不迷失在細節中",
        "堅韌毅力 - 長期主義的堅持",
        "適應能力 - 快速應對市場變化",
        "精力充沛 - 支撐高強度運作"
    ], Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.8), TITLE_COLOR)
    
    add_card(slide5, "5 種能力", [
        "溝通能力 - 建立信任的橋樑",
        "資源整合 - 借力使力放大成果",
        "收集資料 - 用數據分析做決策",
        "敢於冒險 - 在不確定中找機會",
        "承受能力 - 對風險的承受與控管"
    ], Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.8), GREEN_COLOR)
    
    add_card(slide5, "5 大思維", [
        "逆向思維\n從結果倒推原因與路徑",
        "客觀思維\n實事求是，避免主觀偏見",
        "目標思維\n所有行動皆指向清晰結果",
        "換位思考\n站在他人立場與痛點看問題",
        "危機思維\n時刻保持憂患意識與B計畫"
    ], Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.8), HIGHLIGHT_COLOR)

    # ==========================================
    # SLIDE 6: 人際思維
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide6)
    add_slide_header(slide6, "人際思維", "內修心智搞定自己，外順人性搞定關係")
    
    add_card(slide6, "搞定自己 (內修心智)", [
        "和自己和解：接納不完美的自我",
        "不糾結、不內耗、不彆扭",
        "專注於目標，活得更舒適自然"
    ], Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.3), BORDER_COLOR)
    
    add_card(slide6, "搞定別人 (外順利益)", [
        "搞定別人就是搞利益：",
        "宏觀：天下熙熙皆為利來，皆為利往",
        "微觀：成年人不搞無意義社交",
        "社交三看：看背景、看資源、看價值"
    ], Inches(6.9), Inches(1.8), Inches(5.6), Inches(2.3), HIGHLIGHT_COLOR)
    
    # 4 column lower metrics
    metrics = [
        ("見識 (廣度與高度)", ["去過什麼地方", "見過什麼東西", "遇到什麼事情", "擁有多少影響力"], TITLE_COLOR),
        ("背景 (實力與放大)", ["實力 + 能力的結合", "平台是能力的放大器", "遵循合作與互補法則"], BORDER_COLOR),
        ("作用 (資源與對接)", ["明白我是誰的身份", "清晰自己對別人有何用", "懂得把資源對接給誰"], GREEN_COLOR),
        ("態度 (人性與關係)", ["看透七情六慾的本性", "用深度搞好人際關係", "利用規則、順勢做局"], HIGHLIGHT_COLOR)
    ]
    
    for idx, (title, items, color) in enumerate(metrics):
        left_pos = Inches(0.8 + idx * 2.933)
        add_card(slide6, title, items, left_pos, Inches(4.3), Inches(2.733), Inches(2.5), color)

    # ==========================================
    # SLIDE 7: 高手思維
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide7)
    add_slide_header(slide7, "高手思維", "看透人性本質，一切皆為工具")
    
    add_card(slide7, "看透人性", [
        "利益永遠是第一位的",
        "面子上講道義，骨子裡謀利益",
        "明處講義以聚人，暗處取利以謀生",
        "凡事不動真感情，理性看待得失",
        "表面同情弱者，暗地追隨強者",
        "明白錢可以解決生活中 99% 的問題"
    ], Inches(0.8), Inches(1.8), Inches(5.6), Inches(3.5), RED_COLOR)
    
    add_card(slide7, "一切皆工具", [
        "「情」是工具：用以連結人心與情感",
        "「人」是工具：用以資源整合與協作",
        "「錢」是工具：用以解決阻礙與鋪路",
        "求小利者無大成：道德和利益兩頭空",
        "不求小利者有大謀：理性、殘酷而現實",
        "看清遊戲規則，萬物皆為我所用"
    ], Inches(6.9), Inches(1.8), Inches(5.6), Inches(3.5), GREEN_COLOR)
    
    # Bottom callout: Three Qi
    three_qi = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.733), Inches(1.2))
    three_qi.fill.solid()
    three_qi.fill.fore_color.rgb = CARD_BG
    three_qi.line.color.rgb = HIGHLIGHT_COLOR
    three_qi.line.width = Pt(2)
    qi_tf = three_qi.text_frame
    qi_p = qi_tf.paragraphs[0]
    qi_p.alignment = PP_ALIGN.CENTER
    qi_run = qi_p.add_run()
    qi_run.text = "⚡ 城府極深的高手三氣：靜中藏著霸氣，穩中藏著殺氣，極正極邪的匪氣。 ⚡"
    qi_run.font.name = FONT_NAME
    qi_run.font.size = Pt(18)
    qi_run.font.bold = True
    qi_run.font.color.rgb = HIGHLIGHT_COLOR

    # ==========================================
    # SLIDE 8: 處世思維
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_slide_base(prs, slide8)
    add_slide_header(slide8, "處世思維", "處世的三個層次與與人博弈的境界")
    
    levels_3 = [
        ("第一層: 謀事", ["識人大於做事", "溝通大於實幹", "事在人為，人是關鍵"], TITLE_COLOR),
        ("第二層: 謀人", ["平台大於天賦", "方向大於努力", "長度大於速度"], BORDER_COLOR),
        ("第三層: 謀局", ["萬物皆為我所用", "萬物皆不為我所有", "掌握規則，借勢做局"], GREEN_COLOR)
    ]
    
    for idx, (title, items, color) in enumerate(levels_3):
        left_pos = Inches(0.8 + idx * 3.9)
        add_card(slide8, title, items, left_pos, Inches(1.8), Inches(3.7), Inches(2.2), color)
        
    # Bottom split layout
    # Left: Tian Di Ren
    add_card(slide8, "天、地、人 (謀局根本)", [
        "1. 天 (大勢) - 社會發展的必然規律與大趨勢",
        "2. 工具 - 互聯網、人工智慧、雲端計算等生產力",
        "3. 政策 - 方向，研讀大政方針能事半功倍"
    ], Inches(0.8), Inches(4.2), Inches(5.6), Inches(2.4), TITLE_COLOR)
    
    # Right: Two Abilities
    add_card(slide8, "博弈的兩種能力", [
        "1. 化繁為簡\n一眼看透本質，並能抓住事物的主要矛盾",
        "2. 化簡為繁\n找到規律本質後，還能包裝成大眾接納的產品/服務"
    ], Inches(6.9), Inches(4.2), Inches(5.6), Inches(2.4), HIGHLIGHT_COLOR)

    # Save presentation
    output_path = r"d:\Git\true-soul-light\不用版控的\做成投影片的圖\陳老師行銷內容_重繪版.pptx"
    prs.save(output_path)
    print(f"Redrawn presentation saved successfully to: {output_path}")

if __name__ == '__main__':
    main()
