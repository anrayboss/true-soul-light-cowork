import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Color Palette (Elegant Modern White Theme)
BG_COLOR = RGBColor(255, 255, 255)            # Pure White Background
SLIDE_BORDER = RGBColor(226, 232, 240)        # Soft Gray border (Slate-200)

# Card Styling
CARD_BG = RGBColor(248, 250, 252)             # Very soft grey-blue (Slate-50)
CARD_BORDER = RGBColor(226, 232, 240)         # Card boundary line (Slate-200)
ITEM_BG = RGBColor(255, 255, 255)             # Pure White for inner list items
ITEM_BORDER = RGBColor(226, 232, 240)         # Soft border for list items

# Text Colors
TEXT_DARK = RGBColor(15, 23, 42)              # Primary headers (Slate-900)
TEXT_MUTED = RGBColor(71, 85, 105)            # Subtitles (Slate-600)
TEXT_BODY = RGBColor(30, 41, 59)              # Main list content (Slate-800)

# Brand Accents
ACCENT_INDIGO = RGBColor(79, 70, 229)         # Left column accent (Indigo-600)
ACCENT_PURPLE = RGBColor(124, 58, 237)        # Right column accent (Purple-600)

FONT_NAME = 'Microsoft JhengHei'              # Premium Chinese Sans-serif Font

def create_redrawn_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    slide = prs.slides.add_slide(blank_layout)
    
    # 1. Slide Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # 2. Outer Thin Border
    border_margin = Inches(0.1)
    border = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        border_margin, border_margin,
        prs.slide_width - 2 * border_margin,
        prs.slide_height - 2 * border_margin
    )
    border.fill.background()
    border.line.color.rgb = SLIDE_BORDER
    border.line.width = Pt(0.5)
    if border.adjustments:
        border.adjustments[0] = 0.02
        
    # ==========================================
    # LEFT COLUMN: 10 大轉單密碼
    # ==========================================
    left_x = Inches(0.6)
    left_y = Inches(0.6)
    left_w = Inches(7.2)
    left_h = Inches(6.3)
    
    # Left Outer Container Card
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, left_y, left_w, left_h)
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = CARD_BG
    left_card.line.color.rgb = CARD_BORDER
    left_card.line.width = Pt(1)
    if left_card.adjustments:
        left_card.adjustments[0] = 0.03
        
    # Top Accent Line for Left Card
    left_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, left_y, left_w, Inches(0.12))
    left_accent.fill.solid()
    left_accent.fill.fore_color.rgb = ACCENT_INDIGO
    left_accent.line.fill.background()
    
    # Left Titles Text Box
    left_tb = slide.shapes.add_textbox(left_x + Inches(0.3), left_y + Inches(0.3), left_w - Inches(0.6), Inches(1.0))
    left_tf = left_tb.text_frame
    left_tf.word_wrap = True
    left_tf.margin_left = left_tf.margin_top = left_tf.margin_right = left_tf.margin_bottom = 0
    
    p_left_title = left_tf.paragraphs[0]
    p_left_title.space_after = Pt(4)
    run_left_title = p_left_title.add_run()
    run_left_title.text = "10 大轉單密碼"
    run_left_title.font.name = FONT_NAME
    run_left_title.font.size = Pt(26)
    run_left_title.font.bold = True
    run_left_title.font.color.rgb = TEXT_DARK
    
    p_left_sub = left_tf.add_paragraph()
    run_left_sub = p_left_sub.add_run()
    run_left_sub.text = "撬開消費者的心理防線"
    run_left_sub.font.name = FONT_NAME
    run_left_sub.font.size = Pt(13)
    run_left_sub.font.color.rgb = TEXT_MUTED
    
    # 10 Items List Data
    left_items = [
        "獨家賣點",
        "急迫性／稀缺性",
        "談「錢」",
        "掛保證",
        "立權威",
        "成功見證",
        "踩痛點",
        "貼標籤",
        "給利益",
        "隱藏的第二標題：PS"
    ]
    
    # Grid coordinates
    grid_top_start = Inches(1.95)
    row_height = Inches(0.72)
    row_gap = Inches(0.15)
    col_width = Inches(3.15)
    col_gap = Inches(0.3)
    
    for idx, item in enumerate(left_items):
        row = idx % 5
        col = idx // 5
        
        item_x = left_x + Inches(0.3) + col * (col_width + col_gap)
        item_y = grid_top_start + row * (row_height + row_gap)
        
        # Item Rounded Card
        item_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, item_x, item_y, col_width, row_height)
        item_shape.fill.solid()
        item_shape.fill.fore_color.rgb = ITEM_BG
        item_shape.line.color.rgb = ITEM_BORDER
        item_shape.line.width = Pt(1)
        if item_shape.adjustments:
            item_shape.adjustments[0] = 0.08
            
        tf = item_shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        
        # Number Prefix Run
        run_num = p.add_run()
        run_num.text = f"{idx + 1}. "
        run_num.font.name = FONT_NAME
        run_num.font.size = Pt(13)
        run_num.font.bold = True
        run_num.font.color.rgb = ACCENT_INDIGO
        
        # Item Text Run
        run_text = p.add_run()
        run_text.text = item
        run_text.font.name = FONT_NAME
        run_text.font.size = Pt(13)
        run_text.font.bold = True
        run_text.font.color.rgb = TEXT_BODY

    # ==========================================
    # RIGHT COLUMN: 5 大爆文套路
    # ==========================================
    right_x = Inches(8.3)
    right_y = Inches(0.6)
    right_w = Inches(4.4)
    right_h = Inches(6.3)
    
    # Right Outer Container Card
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, right_y, right_w, right_h)
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = CARD_BG
    right_card.line.color.rgb = CARD_BORDER
    right_card.line.width = Pt(1)
    if right_card.adjustments:
        right_card.adjustments[0] = 0.03
        
    # Top Accent Line for Right Card
    right_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, right_y, right_w, Inches(0.12))
    right_accent.fill.solid()
    right_accent.fill.fore_color.rgb = ACCENT_PURPLE
    right_accent.line.fill.background()
    
    # Right Titles Text Box
    right_tb = slide.shapes.add_textbox(right_x + Inches(0.3), right_y + Inches(0.3), right_w - Inches(0.6), Inches(1.0))
    right_tf = right_tb.text_frame
    right_tf.word_wrap = True
    right_tf.margin_left = right_tf.margin_top = right_tf.margin_right = right_tf.margin_bottom = 0
    
    p_right_title = right_tf.paragraphs[0]
    p_right_title.space_after = Pt(4)
    run_right_title = p_right_title.add_run()
    run_right_title.text = "5 大爆文套路"
    run_right_title.font.name = FONT_NAME
    run_right_title.font.size = Pt(26)
    run_right_title.font.bold = True
    run_right_title.font.color.rgb = TEXT_DARK
    
    p_right_sub = right_tf.add_paragraph()
    run_right_sub = p_right_sub.add_run()
    run_right_sub.text = "駕馭流量的爆文套路"
    run_right_sub.font.name = FONT_NAME
    run_right_sub.font.size = Pt(13)
    run_right_sub.font.color.rgb = TEXT_MUTED
    
    # 5 Items List Data
    right_items = [
        "列點文",
        "懶人包",
        "教學文",
        "搬知識",
        "說故事"
    ]
    
    col_width_right = Inches(3.8)
    
    for idx, item in enumerate(right_items):
        item_x = right_x + Inches(0.3)
        item_y = grid_top_start + idx * (row_height + row_gap)
        
        # Item Rounded Card
        item_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, item_x, item_y, col_width_right, row_height)
        item_shape.fill.solid()
        item_shape.fill.fore_color.rgb = ITEM_BG
        item_shape.line.color.rgb = ITEM_BORDER
        item_shape.line.width = Pt(1)
        if item_shape.adjustments:
            item_shape.adjustments[0] = 0.08
            
        tf = item_shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        
        # Number Prefix Run
        run_num = p.add_run()
        run_num.text = f"{idx + 1}. "
        run_num.font.name = FONT_NAME
        run_num.font.size = Pt(13)
        run_num.font.bold = True
        run_num.font.color.rgb = ACCENT_PURPLE
        
        # Item Text Run
        run_text = p.add_run()
        run_text.text = item
        run_text.font.name = FONT_NAME
        run_text.font.size = Pt(13)
        run_text.font.bold = True
        run_text.font.color.rgb = TEXT_BODY

    # Save output PPTX file
    output_folder = r"d:\Git\true-soul-light\不用版控的\做成投影片的圖"
    output_filename = "十大轉單密碼x五大爆文套路_白底投影片.pptx"
    output_path = os.path.join(output_folder, output_filename)
    
    prs.save(output_path)
    print(f"Slide generated successfully at: {output_path}")

if __name__ == "__main__":
    create_redrawn_slide()
