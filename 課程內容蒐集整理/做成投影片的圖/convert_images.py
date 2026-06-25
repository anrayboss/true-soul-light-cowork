import os
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

def natural_sort_key(s):
    # Split by numbers so we can sort numerically
    return [int(text) if text.isdigit() else text.lower() for text in re.split(r'(\d+)', s)]

def convert_images_to_pptx(folder_path, output_pptx_path):
    print(f"Scanning folder: {folder_path}")
    # Get all image files
    supported_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp')
    image_files = [f for f in os.listdir(folder_path) if f.lower().endswith(supported_extensions)]
    
    if not image_files:
        print("No images found in the folder.")
        return
    
    # Sort files naturally
    image_files.sort(key=natural_sort_key)
    print(f"Found {len(image_files)} images to convert.")
    for idx, f in enumerate(image_files):
        print(f"  [{idx+1}] {f}")
        
    # Open first image to get aspect ratio
    first_image_path = os.path.join(folder_path, image_files[0])
    try:
        with Image.open(first_image_path) as img:
            img_width, img_height = img.size
    except Exception as e:
        print(f"Error opening first image: {e}")
        return
        
    prs = Presentation()
    # Define slide dimensions
    # Standard widescreen is 13.333 x 7.5 inches.
    # We will adjust height to match the first image's aspect ratio.
    slide_width = Inches(13.333)
    aspect_ratio = img_height / img_width
    slide_height = Inches(13.333 * aspect_ratio)
    
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    
    blank_layout = prs.slide_layouts[6] # Blank slide layout
    
    for f in image_files:
        image_path = os.path.join(folder_path, f)
        slide = prs.slides.add_slide(blank_layout)
        
        # Add the image covering the full slide
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=slide_width, height=slide_height)
        print(f"Added {f} to slide.")
        
    prs.save(output_pptx_path)
    print(f"Successfully saved to {output_pptx_path}")

if __name__ == '__main__':
    folder = r"d:\Git\true-soul-light\不用版控的\做成投影片的圖"
    output = os.path.join(folder, "做成投影片的圖.pptx")
    convert_images_to_pptx(folder, output)
