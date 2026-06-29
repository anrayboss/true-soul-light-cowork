import os
import sys
import re
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 避免中文編碼輸出問題
sys.stdout.reconfigure(encoding='utf-8')

# 定義配色系統的 RGBColor 物件
COLOR_BG = RGBColor(252, 252, 252)        # #FCFCFC (九陽神功背景色)
COLOR_TEXT_TITLE = RGBColor(8, 51, 68)     # #083344 (Cyan-950 標題文字)
COLOR_ACCENT_BAR = RGBColor(14, 165, 233)  # #0EA5E9 (Sky-500 裝飾條)
COLOR_DIVIDER = RGBColor(226, 232, 240)    # #E2E8F0 (Slate-200 分割線)

COLOR_SKY_50 = RGBColor(240, 249, 255)     # #F0F9FF (主要卡片底色)
COLOR_SKY_100 = RGBColor(224, 242, 254)    # #E0F2FE (主流程/次標底色)
COLOR_SKY_200 = RGBColor(186, 230, 253)    # #BAE6FD (卡片細邊框)
COLOR_SKY_300 = RGBColor(125, 211, 252)    # #7DD3FC (次要流程邊框)
COLOR_SKY_500 = RGBColor(14, 165, 233)     # #0EA5E9 (主流程底色/箭頭)
COLOR_SKY_600 = RGBColor(2, 132, 199)      # #0284C7 (大分類底色)
COLOR_SKY_700 = RGBColor(3, 105, 161)      # #0369A1 (頂部核心底色)

COLOR_SLATE_700 = RGBColor(51, 65, 85)     # #334155 (主文字)
COLOR_WHITE = RGBColor(255, 255, 255)      # 白色

def apply_shape_styling(shape, fill_rgb=None, line_rgb=None, line_width=1.5):
    try:
        if fill_rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
        else:
            shape.fill.background()
            
        if hasattr(shape, 'line') and shape.line:
            if line_rgb:
                shape.line.color.rgb = line_rgb
                shape.line.width = Pt(line_width)
            else:
                if fill_rgb:
                    shape.line.color.rgb = fill_rgb
                else:
                    shape.line.fill.background()
    except Exception:
        pass

def format_text_frame(tf, font_name="微軟正黑體", default_size=12, default_color=COLOR_SLATE_700, bold=False, align=None):
    tf.word_wrap = True
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = align
        p.font.name = font_name
        if p.text.strip():
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(default_size)
                run.font.color.rgb = default_color
                run.font.bold = bold

def add_standard_header(slide, title_text):
    """
    為內容投影片加入九陽風格的標題、Accent Bar 與 Divider Line。
    """
    # 1. 標題文字框
    title_box = slide.shapes.add_textbox(Inches(1.05), Inches(0.60), Inches(11.48), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    p.font.name = "微軟正黑體"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_TITLE
    
    # 2. 天空藍裝飾條 (Accent Bar)
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.65), Inches(0.08), Inches(0.6))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_ACCENT_BAR
    accent_bar.line.fill.background()
    
    # 3. 底部分割線 (Divider Line)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_DIVIDER
    line.line.fill.background()
    
    return title_box

def split_and_add_text_slides(title, elements, slides_list):
    """
    使用「字數折行估算分頁」演算法將文字分割成每頁在畫面上呈現總行數不超過 6 行。
    """
    chunks = []
    current_chunk = []
    current_lines = 0
    
    for el in elements:
        text_len = len(el["text"])
        # 在 24pt/20pt 微軟正黑體與內容寬度下，估算換行折行後的實際行數
        est_lines = (text_len + 21) // 22 if el["type"] == "h3" else (text_len + 27) // 28
        
        # 累加行數，若加上此元素後大於 10 行，則切分到下一頁
        if current_lines + est_lines > 10 and current_chunk:
            chunks.append(current_chunk)
            current_chunk = [el]
            current_lines = est_lines
        else:
            current_chunk.append(el)
            current_lines += est_lines
            
    if current_chunk:
        chunks.append(current_chunk)
        
    for chunk_idx, chunk in enumerate(chunks):
        title_suffix = ""
        if len(chunks) > 1:
            title_suffix = f" ({chunk_idx+1}/{len(chunks)})"
        slides_list.append({
            "title": f"{title}{title_suffix}",
            "type": "text",
            "elements": chunk,
            "table": None
        })

def parse_markdown(filepath):
    """
    解析 Markdown 檔案，將其結構化拆分成 Cover 資訊以及多個章節投影片，並支援二級子章節小表格的左右排版。
    """
    if not os.path.exists(filepath):
        print(f"Error: File not found: {filepath}")
        return None, []
        
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()
        
    content = content.replace('\r\n', '\n')
    
    # 1. 提取第一個 H1 (# ) 作為 Cover
    cover_title = "「真靈光」全球身心靈科技生態系"
    cover_subtitles = []
    
    match_h1 = re.search(r'^#\s*(.*?)(?=\n|$)', content)
    if match_h1:
        cover_title = match_h1.group(1).strip()
        cover_title = re.sub(r'^🌌\s*', '', cover_title)
        
    # 收集副標題
    first_lines = content.split('\n')[:15]
    for line in first_lines:
        line_s = line.strip()
        if not line_s or line_s.startswith('#'):
            continue
        if any(prefix in line_s for prefix in ["呈報", "提案人", "核心理念"]):
            cover_subtitles.append(line_s)
            
    # 2. 按 "## " 切分大章節
    sections_raw = re.split(r'\n(?=##\s+)', content)
    slides_data = []
    
    for sec in sections_raw:
        sec = sec.strip()
        if not sec.startswith('##'):
            continue
            
        lines = sec.split('\n')
        raw_title = lines[0].lstrip('#').strip()
        raw_title = re.sub(r'^(?:📋|🌌|🧱|🔮|🎬|🏢|👑|❓|⚡|💡)\s*', '', raw_title)
        
        content_lines = lines[1:]
        
        # 將 content_lines 拆分成多個 H3 區塊，檢測是否可對 H3 實施左右雙欄重繪
        h3_blocks = []
        current_block = {"h3_title": "", "texts": [], "table_lines": []}
        in_table = False
        h3_count = 0
        
        for line_raw in content_lines:
            line = line_raw.strip()
            if not line:
                continue
                
            if line.startswith('### '):
                h3_count += 1
                if current_block["h3_title"] or current_block["texts"] or current_block["table_lines"]:
                    h3_blocks.append(current_block)
                current_block = {"h3_title": line.split(maxsplit=1)[1].strip(), "texts": [], "table_lines": []}
                in_table = False
                continue
                
            if line.startswith('|'):
                in_table = True
                current_block["table_lines"].append(line_raw)
            else:
                if in_table:
                    in_table = False
                if line.startswith('- ') or line.startswith('* '):
                    leading_spaces = len(line_raw) - len(line_raw.lstrip())
                    level = 1 if leading_spaces >= 2 else 0
                    current_block["texts"].append({"type": "bullet", "text": line[2:].strip(), "level": level})
                elif re.match(r'^\d+\.\s+', line):
                    match = re.match(r'^(\d+)\.\s+(.*)', line)
                    current_block["texts"].append({"type": "bullet", "text": f"{match.group(1)}. {match.group(2).strip()}", "level": 0})
                else:
                    if "目錄" in raw_title:
                        if line.startswith('- ['):
                            continue
                    current_block["texts"].append({"type": "text", "text": line, "level": 0})
                    
        if current_block["h3_title"] or current_block["texts"] or current_block["table_lines"]:
            h3_blocks.append(current_block)
            
        # 決定是否以每一個 H3 作為獨立分頁
        # 條件：章節下有多個 H3，且至少有一組包含小表格（列數 <= 3）
        use_h3_split = False
        if h3_count >= 2 and "四大財務" in raw_title:
            for b in h3_blocks:
                if b["table_lines"]:
                    headers = []
                    for tl in b["table_lines"]:
                        cells = [c.strip() for c in tl.strip().split('|')[1:-1]]
                        if cells and not all(re.match(r'^[-:\s]+$', c) for c in cells):
                            headers = cells
                            break
                    if headers and len(headers) <= 3:
                        use_h3_split = True
                        break
                        
        if use_h3_split:
            for b in h3_blocks:
                h3_title = b["h3_title"]
                full_title = f"{raw_title} - {h3_title}" if h3_title else raw_title
                
                # 解析表格
                parsed_table = None
                if b["table_lines"]:
                    headers = []
                    rows = []
                    for tl in b["table_lines"]:
                        cells = [c.strip() for c in tl.strip().split('|')[1:-1]]
                        if not cells:
                            continue
                        if all(re.match(r'^[-:\s]+$', c) for c in cells):
                            continue
                        if not headers:
                            headers = cells
                        else:
                            rows.append(cells)
                    if headers:
                        parsed_table = {"headers": headers, "rows": rows}
                        
                if parsed_table and len(parsed_table["headers"]) <= 3:
                    # 左右雙欄排版 (左文右表)
                    slides_data.append({
                        "title": full_title,
                        "type": "text_table",
                        "elements": b["texts"],
                        "table": parsed_table
                    })
                else:
                    # 文字與大表格分頁
                    if b["texts"]:
                        split_and_add_text_slides(full_title, b["texts"], slides_data)
                    if parsed_table:
                        slides_data.append({
                            "title": full_title,
                            "type": "table",
                            "elements": [],
                            "table": parsed_table
                        })
            continue
            
        # 若不走二級 H3 拆分，則維持整區塊的大分頁解析
        text_elements = []
        table_lines = []
        in_table = False
        
        for line_raw in content_lines:
            line = line_raw.strip()
            if not line:
                continue
            if line.startswith('|'):
                in_table = True
                table_lines.append(line_raw)
                continue
            else:
                in_table = False
                
            if line.startswith('### '):
                text_elements.append({"type": "h3", "text": line.split(maxsplit=1)[1].strip(), "level": 0})
            elif line.startswith('- ') or line.startswith('* '):
                leading_spaces = len(line_raw) - len(line_raw.lstrip())
                level = 1 if leading_spaces >= 2 else 0
                text_elements.append({"type": "bullet", "text": line[2:].strip(), "level": level})
            elif re.match(r'^\d+\.\s+', line):
                match = re.match(r'^(\d+)\.\s+(.*)', line)
                text_elements.append({"type": "bullet", "text": f"{match.group(1)}. {match.group(2).strip()}", "level": 0})
            else:
                if "目錄" in raw_title:
                    if line.startswith('- ['):
                        continue
                text_elements.append({"type": "text", "text": line, "level": 0})
                
        parsed_table = None
        if table_lines:
            headers = []
            rows = []
            for tl in table_lines:
                cells = [c.strip() for c in tl.strip().split('|')[1:-1]]
                if not cells:
                    continue
                if all(re.match(r'^[-:\s]+$', c) for c in cells):
                    continue
                if not headers:
                    headers = cells
                else:
                    rows.append(cells)
            if headers:
                parsed_table = {"headers": headers, "rows": rows}
                
        # 智慧分頁
        if text_elements and parsed_table:
            split_and_add_text_slides(raw_title, text_elements, slides_data)
            slides_data.append({
                "title": raw_title,
                "type": "table",
                "elements": [],
                "table": parsed_table
            })
        elif parsed_table:
            slides_data.append({
                "title": raw_title,
                "type": "table",
                "elements": [],
                "table": parsed_table
            })
        elif text_elements:
            split_and_add_text_slides(raw_title, text_elements, slides_data)
            
    return {"title": cover_title, "subtitles": cover_subtitles}, slides_data

def build_pptx(cover_data, slides_list, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ==================== 1. 建立封面頁 ====================
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(3.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = cover_data["title"]
    p.font.name = '微軟正黑體'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_TITLE
    p.alignment = PP_ALIGN.CENTER
    
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.3), Inches(4.333), Inches(0.03))
    divider.fill.solid()
    divider.fill.fore_color.rgb = COLOR_ACCENT_BAR
    divider.line.fill.background()
    
    if cover_data["subtitles"]:
        sub_text = " | ".join(cover_data["subtitles"])
        p2 = tf.add_paragraph()
        p2.text = sub_text
        p2.font.name = '微軟正黑體'
        p2.font.size = Pt(18)
        p2.font.color.rgb = COLOR_SKY_700
        p2.space_before = Pt(40)
        p2.alignment = PP_ALIGN.CENTER
        
    # ==================== 2. 建立內容頁 ====================
    for slide_data in slides_list:
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLOR_BG
        
        # 建立標題
        add_standard_header(slide, slide_data["title"])
        
        if slide_data["type"] == "text":
            # 建立文字框 (L=0.8", T=1.85", W=11.73", H=4.77")
            tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(11.73), Inches(4.77))
            tf = tx_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            
            for idx, el in enumerate(slide_data["elements"]):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = el["text"]
                p.font.name = "微軟正黑體"
                
                # 根據層級與類型設定樣式 (最小字體 20pt，H3 為 24pt)
                if el["type"] == "h3":
                    p.font.size = Pt(24)
                    p.font.bold = True
                    p.font.color.rgb = COLOR_TEXT_TITLE
                    p.space_before = Pt(16)
                elif el["type"] == "bullet":
                    p.font.size = Pt(20)
                    p.font.bold = False
                    p.font.color.rgb = COLOR_SLATE_700
                    p.level = el["level"]
                    p.space_before = Pt(12)
                else:
                    p.font.size = Pt(20)
                    p.font.bold = False
                    p.font.color.rgb = COLOR_SLATE_700
                    p.space_before = Pt(12)
                    
        elif slide_data["type"] == "text_table":
            # 左右雙欄排版 (左邊文字說明 20pt / 右邊表格對照)
            # 1. 左欄文字說明 (L=0.8", T=1.85", W=4.2", H=4.5")
            tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(4.2), Inches(4.5))
            tf = tx_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            
            for idx, el in enumerate(slide_data["elements"]):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = el["text"]
                p.font.name = "微軟正黑體"
                
                if el["type"] == "h3":
                    p.font.size = Pt(24)
                    p.font.bold = True
                    p.font.color.rgb = COLOR_TEXT_TITLE
                    p.space_before = Pt(16)
                elif el["type"] == "bullet":
                    p.font.size = Pt(20)
                    p.font.bold = False
                    p.font.color.rgb = COLOR_SLATE_700
                    p.level = el["level"]
                    p.space_before = Pt(12)
                else:
                    p.font.size = Pt(20)
                    p.font.bold = False
                    p.font.color.rgb = COLOR_SLATE_700
                    p.space_before = Pt(12)
                    
            # 2. 右欄表格 (L=5.4", T=1.85", W=7.1", H=4.5")
            table_data = slide_data["table"]
            rows_count = len(table_data["rows"]) + 1
            cols_count = len(table_data["headers"])
            table_h = Inches(min(4.8, 0.5 * rows_count))
            
            table_shape = slide.shapes.add_table(rows_count, cols_count, Inches(5.4), Inches(1.85), Inches(7.1), table_h)
            table = table_shape.table
            
            # 格式化表頭
            for col_i, header in enumerate(table_data["headers"]):
                cell = table.cell(0, col_i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_SKY_600
                cell.text = header
                format_text_frame(cell.text_frame, default_size=14, default_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
                cell.margin_left = cell.margin_right = Inches(0.05)
                cell.margin_top = cell.margin_bottom = Inches(0.10)
                
            # 格式化內容
            for row_i, row in enumerate(table_data["rows"]):
                for col_i, text in enumerate(row):
                    cell = table.cell(row_i + 1, col_i)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_SKY_50
                    cell.text = text
                    format_text_frame(cell.text_frame, default_size=13, default_color=COLOR_SLATE_700, bold=False, align=PP_ALIGN.CENTER)
                    cell.margin_left = cell.margin_right = Inches(0.05)
                    cell.margin_top = cell.margin_bottom = Inches(0.10)
                    
        elif slide_data["type"] == "table":
            table_data = slide_data["table"]
            rows_count = len(table_data["rows"]) + 1
            cols_count = len(table_data["headers"])
            
            table_w = Inches(11.73)
            table_l = Inches(0.8)
            table_t = Inches(1.85)
            table_h = Inches(min(4.8, 0.55 * rows_count))
            
            table_shape = slide.shapes.add_table(rows_count, cols_count, table_l, table_t, table_w, table_h)
            table = table_shape.table
            
            for col_i, header in enumerate(table_data["headers"]):
                cell = table.cell(0, col_i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_SKY_600
                cell.text = header
                format_text_frame(cell.text_frame, default_size=16, default_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
                cell.margin_left = cell.margin_right = Inches(0.05)
                cell.margin_top = cell.margin_bottom = Inches(0.12)
                
            for row_i, row in enumerate(table_data["rows"]):
                for col_i, text in enumerate(row):
                    cell = table.cell(row_i + 1, col_i)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_SKY_50
                    cell.text = text
                    format_text_frame(cell.text_frame, default_size=14, default_color=COLOR_SLATE_700, bold=False, align=PP_ALIGN.CENTER)
                    cell.margin_left = cell.margin_right = Inches(0.05)
                    cell.margin_top = cell.margin_bottom = Inches(0.12)
                    
    try:
        prs.save(output_path)
        print(f"Presentation built successfully: {output_path}")
    except PermissionError:
        alternative_path = output_path.replace(".pptx", "_v2.pptx")
        print(f"\n[WARNING] Permission Denied: Could not write to {output_path}.")
        print("Please close the file in PowerPoint if it is open.")
        print(f"Saving to alternative path instead: {alternative_path}")
        prs.save(alternative_path)
        print(f"Presentation built successfully: {alternative_path}")

def main():
    md_path = r"真靈光企劃書所有資訊(線上共編處)\真靈光企劃書_正式提案版完整版本.md"
    output_pptx = r"真靈光企劃書所有資訊(線上共編處)\真靈光企劃書_正式提案版完整版本.pptx"
    
    print("Parsing proposal markdown with smart layouts...")
    cover_data, slides_list = parse_markdown(md_path)
    
    if not slides_list:
        print("Error: No slides parsed!")
        sys.exit(1)
        
    print(f"Parsed {len(slides_list)} slides (including H3-split and line-wrapped slides).")
    
    print("Building PPTX with Nine-Yang styling...")
    build_pptx(cover_data, slides_list, output_pptx)

if __name__ == "__main__":
    main()
