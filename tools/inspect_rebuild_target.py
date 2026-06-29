import sys
import io
import os
from pptx import Presentation

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def inspect_pptx(path, name):
    print(f"\n=========================================")
    print(f"Inspecting: {name} ({path})")
    if not os.path.exists(path):
        print("File not found!")
        return
    prs = Presentation(path)
    print(f"Total slides: {len(prs.slides)}")
    
    # Print presentation size
    w = prs.slide_width.inches
    h = prs.slide_height.inches
    print(f"Slide Size: Width={w:.3f}\", Height={h:.3f}\"")
    
    for i, slide in enumerate(prs.slides): # inspect first 10 slides for details
        if i >= 10:
            break
        print(f"\n--- Slide {i+1} ---")
        # Check background fill
        bg = slide.background
        if bg and bg.fill:
            print(f"  Background type: {bg.fill.type}")
            if bg.fill.type == 1: # Solid
                print(f"  Background Color: {bg.fill.fore_color.rgb if bg.fill.fore_color else 'None'}")
        
        for shape in slide.shapes:
            shape_type = shape.shape_type
            shape_name = shape.name
            left = shape.left.inches if shape.left else 0
            top = shape.top.inches if shape.top else 0
            width = shape.width.inches if shape.width else 0
            height = shape.height.inches if shape.height else 0
            text_preview = ""
            if shape.has_text_frame:
                texts = []
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
                text_preview = " | Texts: " + " / ".join(texts)[:200]
            print(f"  Shape: '{shape_name}', Type: {shape_type}, Pos: L={left:.2f}\", T={top:.2f}\", W={width:.2f}\", H={height:.2f}\"{text_preview}")

print("=== REBUILD TARGET INSTRUCTION ===")
inspect_pptx(r"真靈光企劃書所有資訊(線上共編處)\重製.pptx", "Rebuild Target")

print("\n=== STYLE REFERENCE INSTRUCTION ===")
inspect_pptx(r"最終輸出用檔案\舊版、已使用封存\行銷界九陽神功_完整講稿_總監編修#4.pptx", "Style Reference")
