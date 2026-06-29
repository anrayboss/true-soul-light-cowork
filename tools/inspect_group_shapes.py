import os
import sys
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

def inspect_shapes_recursive(shape, depth, f):
    indent = "  " * depth
    shape_type = shape.shape_type
    shape_name = shape.name
    left = shape.left.inches if shape.left else 0
    top = shape.top.inches if shape.top else 0
    width = shape.width.inches if shape.width else 0
    height = shape.height.inches if shape.height else 0
    
    text_preview = ""
    if shape.has_text_frame:
        texts = []
        for p in shape.text_frame.paragraphs:
            if p.text.strip():
                texts.append(p.text.strip())
        text_preview = " | Texts: " + " / ".join(texts)
        
    f.write(f"{indent}Shape: '{shape_name}', Type: {shape_type}, Pos: L={left:.2f}\", T={top:.2f}\", W={width:.2f}\", H={height:.2f}\"{text_preview}\n")
    
    # Check if group shape
    if shape_type == 6 or hasattr(shape, "shapes"): # GROUP is 6
        try:
            for subshape in shape.shapes:
                inspect_shapes_recursive(subshape, depth + 1, f)
        except Exception as e:
            f.write(f"{indent}  [Error recursing group: {str(e)}]\n")

def inspect_pptx_detail(path, output_log):
    with open(output_log, 'w', encoding='utf-8') as f:
        f.write(f"Detailed Inspection of: {path}\n")
        if not os.path.exists(path):
            f.write("File not found!\n")
            return
        
        try:
            prs = Presentation(path)
            f.write(f"Total slides: {len(prs.slides)}\n")
            f.write(f"Slide Size: Width={prs.slide_width.inches:.3f}\", Height={prs.slide_height.inches:.3f}\"\n")
            
            for i, slide in enumerate(prs.slides):
                f.write(f"\n--- Slide {i+1} ---\n")
                
                # Check background
                bg = slide.background
                if bg and bg.fill:
                    f.write(f"  Background Fill Type: {bg.fill.type}\n")
                    
                for shape in slide.shapes:
                    inspect_shapes_recursive(shape, 1, f)
                    
            print(f"Detail inspection complete. Saved to {output_log}")
        except Exception as e:
            f.write(f"Error: {str(e)}\n")
            print(f"Error: {str(e)}")

inspect_pptx_detail(r"真靈光企劃書所有資訊(線上共編處)\真靈光企劃書_給 AI 看的.pptx", "tools/ai_seen_pptx_detail.txt")
