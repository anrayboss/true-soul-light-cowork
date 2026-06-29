import sys
import os
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

def verify_slide_shapes(shape, slide_idx):
    shape_type = shape.shape_type
    shape_name = shape.name
    
    # 遞迴處理 Group
    if shape_type == 6 or hasattr(shape, "shapes"):
        for subshape in shape.shapes:
            verify_slide_shapes(subshape, slide_idx)
        return
        
    text = ""
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
        
        # 檢查段落與 Run 的字型
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                # 檢查字型是否為微軟正黑體
                font_name = run.font.name
                size_pt = run.font.size.pt if run.font.size else None
                try:
                    rgb = run.font.color.rgb if run.font.color else None
                except Exception:
                    rgb = None
                
                # 如果有新細明體或 PMingLiU 則警告
                if font_name and ("PMingLiU" in font_name or "新細明體" in font_name):
                    print(f"[WARNING] Slide {slide_idx+1}: Shape '{shape_name}' contains illegal font '{font_name}' in run '{run.text[:20]}'")
                
                # 檢查是否為 None
                if not font_name:
                    # 這通常代表繼承自簡報母片，在我們的代碼中，我們強烈希望全部顯式指定為微軟正黑體
                    # 讓我們看是否有繼承問題
                    pass

    # 檢查填充色與邊框
    try:
        if shape.fill and shape.fill.type == 1:
            rgb_val = shape.fill.fore_color.rgb
            # 印出帶文字的卡片顏色
            if text:
                print(f"Slide {slide_idx+1}: Card '{text[:15]}' -> Fill: {rgb_val}, Line: {shape.line.color.rgb if shape.line else 'None'}")
    except Exception:
        pass

def main():
    path = r"真靈光企劃書所有資訊(線上共編處)\真靈光企劃書_正式提案版完整版本_v2.pptx"
    if not os.path.exists(path):
        path = r"真靈光企劃書所有資訊(線上共編處)\真靈光企劃書_正式提案版完整版本.pptx"
    print(f"Verifying reconstructed presentation: {path}")
    prs = Presentation(path)
    for idx, slide in enumerate(prs.slides):
        print(f"\n--- Verifying Slide {idx+1} ---")
        # 檢查背景色
        bg = slide.background
        if bg and bg.fill and bg.fill.type == 1:
            print(f"Background Color: {bg.fill.fore_color.rgb}")
        else:
            print(f"Background Type: {bg.fill.type if bg and bg.fill else 'None'}")
            
        for shape in slide.shapes:
            verify_slide_shapes(shape, idx)

if __name__ == "__main__":
    main()
