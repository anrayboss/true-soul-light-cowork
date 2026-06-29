import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 避免中文編碼輸出問題
sys.stdout.reconfigure(encoding='utf-8')

# 定義配色系統的 RGBColor 物件
COLOR_BG = RGBColor(252, 252, 252)        # #FCFCFC (九陽神功背景色)
COLOR_TEXT_TITLE = RGBColor(8, 51, 68)     # #083344 (Cyan-950 標題文字)
COLOR_ACCENT_BAR = RGBColor(14, 165, 233)  # #0EA5E9 (Sky-500 裝飾條)
COLOR_DIVIDER = RGBColor(226, 232, 240)    # #E2E8F0 (Slate-200 分割線)

COLOR_SKY_50 = RGBColor(240, 249, 255)     # #F0F9FF (主要卡片底色)
COLOR_SKY_100 = RGBColor(224, 242, 254)    # #E0F2FE (主流程/次標底色)
COLOR_SKY_200 = RGBColor(186, 230, 253)    # #BAE6FD (卡片細邊框)
COLOR_SKY_300 = RGBColor(125, 211, 252)    # #7DD3FC (次要流程邊框)
COLOR_SKY_500 = RGBColor(14, 165, 233)     # #0EA5E9 (主流程底色/箭頭)
COLOR_SKY_600 = RGBColor(2, 132, 199)      # #0284C7 (大分類底色)
COLOR_SKY_700 = RGBColor(3, 105, 161)      # #0369A1 (頂部核心底色)

COLOR_SLATE_700 = RGBColor(51, 65, 85)     # #334155 (主文字)
COLOR_WHITE = RGBColor(255, 255, 255)      # 白色

def format_text_frame(tf, font_name="微軟正黑體", default_size=12, default_color=COLOR_SLATE_700, bold=False, align=None):
    tf.word_wrap = True
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = align
        p.font.name = font_name
        if p.text.strip():
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(default_size)
                run.font.color.rgb = default_color
                run.font.bold = bold

def apply_shape_styling(shape, fill_rgb=None, line_rgb=None, line_width=1.5):
    try:
        if fill_rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
        else:
            shape.fill.background()
            
        if hasattr(shape, 'line') and shape.line:
            if line_rgb:
                shape.line.color.rgb = line_rgb
                shape.line.width = Pt(line_width)
            else:
                if fill_rgb:
                    shape.line.color.rgb = fill_rgb
                else:
                    shape.line.fill.background()
    except Exception:
        pass

def clear_slide_content(slide, title_shape):
    """
    清空簡報內除了標題、Accent Bar 和 Divider Line 以外的所有舊 shape。
    """
    to_delete = []
    for shape in slide.shapes:
        if shape == title_shape:
            continue
        left = shape.left.inches if shape.left else 0
        top = shape.top.inches if shape.top else 0
        # 保留左側的 Accent Bar 和下面的 Divider Line (通常在 L=0.8)
        if left == 0.8 and (top == 0.65 or top == 1.45):
            continue
        to_delete.append(shape)
        
    for shape in to_delete:
        try:
            shape.element.getparent().remove(shape.element)
        except Exception:
            pass

def draw_line(slide, x1, y1, x2, y2, color, thickness=1.5):
    """
    以細長方形替代直線，在 python-pptx 中最為穩定，且便於控制顏色。
    """
    if abs(y1 - y2) < 0.01:
        # 水平線
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1, y1, abs(x2 - x1), Inches(thickness / 72.0))
    else:
        # 垂直線
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1, y1, Inches(thickness / 72.0), abs(y2 - y1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def draw_arrow(slide, direction, left, top, width, height, color):
    shape_type = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    arrow = slide.shapes.add_shape(shape_type, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow

def draw_card(slide, text, left, top, width, height, fill_color, border_color=None, border_width=1.5, font_size=12, font_color=COLOR_SLATE_700, bold=False, align=PP_ALIGN.CENTER):
    """
    繪製一個圓角矩形卡片，先填入文字，再設定格式以防格式被覆寫。
    """
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    apply_shape_styling(card, fill_rgb=fill_color, line_rgb=border_color, line_width=border_width)
    
    tf = card.text_frame
    tf.text = text
    # 文字填入後，重新對文字框內部的 runs 進行格式化
    format_text_frame(tf, default_size=font_size, default_color=font_color, bold=bold, align=align)
    return card

# ==================== 重繪各頁面函數 ====================

def recreate_slide_5(slide, title_shape):
    """
    重繪 Slide 5 (身心靈流程)：三欄 (流量 ➔ 變現 ➔ 倍增)
    """
    clear_slide_content(slide, title_shape)
    
    # 欄寬度與間距
    col_w = Inches(3.4)
    gap = Inches(0.55)
    
    # 三欄的標題與內容
    cols = [
        {
            "title": "流量",
            "text": "天：網站、自媒體、新聞、SEO\n\n地：店面\n\n人：代理商",
            "left": Inches(0.8)
        },
        {
            "title": "變現",
            "text": "諮詢或一對一 (靈氣、諮詢)\n\n課程 (身心靈課程)\n\n產品 (健康食品、儀器)",
            "left": Inches(0.8 + 3.4 + 0.55)
        },
        {
            "title": "倍增",
            "text": "代理商\n\n加盟店",
            "left": Inches(0.8 + 2 * (3.4 + 0.55))
        }
    ]
    
    for col in cols:
        L = col["left"]
        # 1. 畫欄位標題 (Sky-500 滿版)
        draw_card(slide, col["title"], L, Inches(2.0), col_w, Inches(0.6), 
                  fill_color=COLOR_SKY_500, font_size=16, font_color=COLOR_WHITE, bold=True)
                  
        # 2. 畫欄位內容 (Sky-50 背景、Sky-200 邊框)
        draw_card(slide, col["text"], L, Inches(2.75), col_w, Inches(3.6), 
                  fill_color=COLOR_SKY_50, border_color=COLOR_SKY_200, border_width=1.5,
                  font_size=14, font_color=COLOR_SLATE_700, bold=False, align=PP_ALIGN.LEFT)
                  
    # 3. 畫流程箭頭 (Sky-500 滿版)
    draw_arrow(slide, "right", Inches(4.35), Inches(2.05), Inches(0.35), Inches(0.5), COLOR_SKY_500)
    draw_arrow(slide, "right", Inches(8.30), Inches(2.05), Inches(0.35), Inches(0.5), COLOR_SKY_500)

def recreate_slide_7(slide, title_shape):
    """
    重繪 Slide 7 (身心靈一站式服務培訓系統)：4 欄流向圖
    """
    clear_slide_content(slide, title_shape)
    
    # 頂部核心平台
    draw_card(slide, "身心靈整合平台", Inches(4.41), Inches(1.65), Inches(4.5), Inches(0.65), 
              fill_color=COLOR_SKY_700, font_size=15, font_color=COLOR_WHITE, bold=True)
              
    draw_card(slide, "整合一致性", Inches(4.91), Inches(2.4), Inches(3.5), Inches(0.4), 
              fill_color=COLOR_SKY_100, border_color=COLOR_SKY_300, font_size=11, font_color=COLOR_TEXT_TITLE, bold=True)
              
    # 4 欄的座標與定義
    col_w = Inches(2.65)
    gap = Inches(0.38)
    
    channels = [
        {
            "left": Inches(0.8),
            "step1": "宣傳",
            "step2": "流量",
            "step3": "行銷與業務培訓",
            "step4": "流量變現"
        },
        {
            "left": Inches(0.8 + 2.65 + 0.38),
            "step1": "量測",
            "step2": "可數據化",
            "step3": "AI身心靈檢測",
            "step4": "身心靈量測"
        },
        {
            "left": Inches(0.8 + 2 * (2.65 + 0.38)),
            "step1": "諮詢 / 課程",
            "step2": "獨特技術諮詢",
            "step3": "身心靈老師培訓",
            "step4": "可複製靈性提升"
        },
        {
            "left": Inches(0.8 + 3 * (2.65 + 0.38)),
            "step1": "產品",
            "step2": "建議顧客所需",
            "step3": "產品培訓",
            "step4": "獨特產品"
        }
    ]
    
    # 樹狀連接線：從頂部平台向下分支
    # 畫一條橫跨 4 欄中心點的水平線
    draw_line(slide, Inches(2.125), Inches(2.95), Inches(11.215), Inches(2.95), COLOR_SKY_500, thickness=2.0)
    # 從平台連下來
    draw_line(slide, Inches(6.666), Inches(2.8), Inches(6.666), Inches(2.95), COLOR_SKY_500, thickness=2.0)
    
    for ch in channels:
        L = ch["left"]
        center_x = L + col_w / 2
        
        # 連接橫線到第 1 步
        draw_line(slide, center_x, Inches(2.95), center_x, Inches(3.1), COLOR_SKY_500, thickness=2.0)
        
        # 1. 階段名稱 (Sky-500 滿版)
        draw_card(slide, ch["step1"], L, Inches(3.1), col_w, Inches(0.5), 
                  fill_color=COLOR_SKY_500, font_size=13, font_color=COLOR_WHITE, bold=True)
                  
        # 箭頭 1 ➔ 2
        draw_arrow(slide, "down", center_x - Inches(0.12), Inches(3.65), Inches(0.24), Inches(0.2), COLOR_SKY_500)
        
        # 2. 客戶流程說明 (Sky-50 背景、Sky-200 邊框)
        draw_card(slide, ch["step2"], L, Inches(3.9), col_w, Inches(0.55), 
                  fill_color=COLOR_SKY_50, border_color=COLOR_SKY_200, font_size=11, font_color=COLOR_SLATE_700, bold=True)
                  
        # 箭頭 2 ➔ 3
        draw_arrow(slide, "down", center_x - Inches(0.12), Inches(4.5), Inches(0.24), Inches(0.2), COLOR_SKY_500)
        
        # 3. 培訓系統核心 (Sky-100 背景、Sky-300 邊框)
        draw_card(slide, ch["step3"], L, Inches(4.75), col_w, Inches(0.65), 
                  fill_color=COLOR_SKY_100, border_color=COLOR_SKY_300, font_size=12, font_color=COLOR_TEXT_TITLE, bold=True)
                  
        # 箭頭 3 ➔ 4
        draw_arrow(slide, "down", center_x - Inches(0.12), Inches(5.45), Inches(0.24), Inches(0.25), COLOR_SKY_500)
        
        # 4. 細部成果卡片 (Sky-50 背景、Sky-200 邊框)
        draw_card(slide, ch["step4"], L, Inches(5.75), col_w, Inches(0.55), 
                  fill_color=COLOR_SKY_50, border_color=COLOR_SKY_200, font_size=11, font_color=COLOR_SLATE_700, bold=False)

def recreate_slide_10(slide, title_shape):
    """
    重繪 Slide 10 (能量 / 企業五量關係圖)
    """
    clear_slide_content(slide, title_shape)
    
    # 第一層 (CEO 與兩側關聯)
    draw_card(slide, "CEO", Inches(5.36), Inches(1.8), Inches(2.6), Inches(0.65), 
              fill_color=COLOR_SKY_700, font_size=14, font_color=COLOR_WHITE, bold=True)
              
    draw_card(slide, "真XX集團", Inches(1.8), Inches(1.8), Inches(2.6), Inches(0.65), 
              fill_color=COLOR_SKY_50, border_color=COLOR_SKY_200, font_size=12, font_color=COLOR_TEXT_TITLE, bold=True)
              
    draw_card(slide, "陳老師\n(最高隱藏顧問)", Inches(8.9), Inches(1.8), Inches(2.6), Inches(0.65), 
              fill_color=COLOR_SKY_50, border_color=COLOR_SKY_200, font_size=12, font_color=COLOR_TEXT_TITLE, bold=True)
              
    # 連接線 (水平)
    draw_line(slide, Inches(4.4), Inches(2.12), Inches(5.36), Inches(2.12), COLOR_SKY_500, thickness=2.0)
    draw_line(slide, Inches(7.96), Inches(2.12), Inches(8.9), Inches(2.12), COLOR_SKY_500, thickness=2.0)
    
    # 5 欄的分支配置
    col_w = Inches(2.1)
    gap = Inches(0.3)
    
    branches = [
        {
            "left": Inches(0.8),
            "biz": "TIKTOK公會\n(流量來源)",
            "val": "流量"
        },
        {
            "left": Inches(0.8 + 2.4),
            "biz": "真XX身心靈平台\n(第一品牌)",
            "val": "存量"
        },
        {
            "left": Inches(0.8 + 2 * 2.4),
            "biz": "真XX社群媒體\n(真的facebook)",
            "val": "大量"
        },
        {
            "left": Inches(0.8 + 3 * 2.4),
            "biz": "產品製造\n(量子能量&訊息製造廠)",
            "val": "能量"
        },
        {
            "left": Inches(0.8 + 4 * 2.4),
            "biz": "健康食品\n(專利)",
            "val": "使用量"
        }
    ]
    
    # 樹狀連接線：從 CEO 往下連至 5 分支
    # 水平大線
    draw_line(slide, Inches(1.85), Inches(3.0), Inches(11.45), Inches(3.0), COLOR_SKY_500, thickness=2.0)
    # 從 CEO 往下連
    draw_line(slide, Inches(6.66), Inches(2.45), Inches(6.66), Inches(3.0), COLOR_SKY_500, thickness=2.0)
    
    for br in branches:
        L = br["left"]
        center_x = L + col_w / 2
        
        # 從水平大線連到分支卡片
        draw_line(slide, center_x, Inches(3.0), center_x, Inches(3.4), COLOR_SKY_500, thickness=2.0)
        
        # 1. 第二層：分公司/業務 (Sky-100 背景、Sky-300 邊框)
        draw_card(slide, br["biz"], L, Inches(3.4), col_w, Inches(1.2), 
                  fill_color=COLOR_SKY_100, border_color=COLOR_SKY_300, font_size=11, font_color=COLOR_TEXT_TITLE, bold=True)
                  
        # 垂直箭頭連到第三層
        draw_arrow(slide, "down", center_x - Inches(0.12), Inches(4.75), Inches(0.24), Inches(0.45), COLOR_SKY_500)
        
        # 2. 第三層：量指標 (Sky-500 滿版，白字)
        draw_card(slide, br["val"], L, Inches(5.3), col_w, Inches(0.65), 
                  fill_color=COLOR_SKY_500, font_size=14, font_color=COLOR_WHITE, bold=True)

def recreate_slide_11(slide, title_shape):
    """
    重繪 Slide 11 (企業架構組織圖)
    """
    clear_slide_content(slide, title_shape)
    
    # Level 1: 負責人 ➔ CEO
    draw_card(slide, "負責人", Inches(5.36), Inches(1.5), Inches(2.6), Inches(0.55), 
              fill_color=COLOR_SKY_700, font_size=13, font_color=COLOR_WHITE, bold=True)
              
    draw_card(slide, "CEO", Inches(5.36), Inches(2.45), Inches(2.6), Inches(0.55), 
              fill_color=COLOR_SKY_700, font_size=13, font_color=COLOR_WHITE, bold=True)
              
    # 負責人 ➔ CEO 箭頭
    draw_arrow(slide, "down", Inches(6.54), Inches(2.1), Inches(0.24), Inches(0.3), COLOR_SKY_500)
    
    # 部門 L2 配置 (教育、行銷、行政)
    col_w = Inches(3.2)
    gap = Inches(0.76)
    
    # 畫一條橫線用以分支出三大部門
    draw_line(slide, Inches(2.4), Inches(3.25), Inches(10.32), Inches(3.25), COLOR_SKY_500, thickness=2.0)
    # 從 CEO 指向橫線
    draw_line(slide, Inches(6.66), Inches(3.0), Inches(6.66), Inches(3.25), COLOR_SKY_500, thickness=2.0)
    
    # 三個主要部門卡片 (Sky-500 滿版)
    # 1. 教育 (培訓總監)
    draw_card(slide, "教育\n(培訓總監)", Inches(0.8), Inches(3.45), col_w, Inches(0.65), 
              fill_color=COLOR_SKY_500, font_size=13, font_color=COLOR_WHITE, bold=True)
    draw_line(slide, Inches(2.4), Inches(3.25), Inches(2.4), Inches(3.45), COLOR_SKY_500, thickness=2.0)
    
    # 2. 行銷 (行銷總監)
    draw_card(slide, "行銷\n(行銷總監)", Inches(4.76), Inches(3.45), col_w, Inches(0.65), 
              fill_color=COLOR_SKY_500, font_size=13, font_color=COLOR_WHITE, bold=True)
    draw_line(slide, Inches(6.36), Inches(3.25), Inches(6.36), Inches(3.45), COLOR_SKY_500, thickness=2.0)
    
    # 3. 行政 (行政總監)
    draw_card(slide, "行政\n(行政總監)", Inches(8.72), Inches(3.45), col_w, Inches(0.65), 
              fill_color=COLOR_SKY_500, font_size=13, font_color=COLOR_WHITE, bold=True)
    draw_line(slide, Inches(10.32), Inches(3.25), Inches(10.32), Inches(3.45), COLOR_SKY_500, thickness=2.0)
    
    # ==================== 子層級繪製 (Level 3 & 4) ====================
    
    # A. 教育組的子組 (業務培訓、產品培訓) - X: 0.8" ~ 4.0"
    draw_arrow(slide, "down", Inches(2.28), Inches(4.15), Inches(0.24), Inches(0.35), COLOR_SKY_500)
    draw_card(slide, "業務培訓", Inches(0.8), Inches(4.55), Inches(1.4), Inches(0.65), 
              fill_color=COLOR_SKY_100, border_color=COLOR_SKY_300, font_size=11, font_color=COLOR_TEXT_TITLE, bold=True)
    draw_card(slide, "產品培訓", Inches(2.6), Inches(4.55), Inches(1.4), Inches(0.65), 
              fill_color=COLOR_SKY_100, border_color=COLOR_SKY_300, font_size=11, font_color=COLOR_TEXT_TITLE, bold=True)
              
    # B. 行政組的子組 (財務、客服、技術) - X: 8.72" ~ 11.92"
    draw_arrow(slide, "down", Inches(10.2), Inches(4.15), Inches(0.24), Inches(0.35), COLOR_SKY_500)
    draw_card(slide, "財務", Inches(8.72), Inches(4.55), Inches(0.95), Inches(0.65), 
              fill_color=COLOR_SKY_100, border_color=COLOR_SKY_300, font_size=10, font_color=COLOR_TEXT_TITLE, bold=True)
    draw_card(slide, "客服", Inches(9.82), Inches(4.55), Inches(0.95), Inches(0.65), 
              fill_color=COLOR_SKY_100, border_color=COLOR_SKY_300, font_size=10, font_color=COLOR_TEXT_TITLE, bold=True)
    draw_card(slide, "技術", Inches(10.92), Inches(4.55), Inches(0.95), Inches(0.65), 
              fill_color=COLOR_SKY_100, border_color=COLOR_SKY_300, font_size=10, font_color=COLOR_TEXT_TITLE, bold=True)
              
    # C. 行銷組的子組 (網路、店、業務) - X: 4.76" ~ 7.96"
    draw_arrow(slide, "down", Inches(6.24), Inches(4.15), Inches(0.24), Inches(0.3), COLOR_SKY_500)
    draw_card(slide, "網路", Inches(4.6), Inches(4.5), Inches(1.2), Inches(0.55), 
              fill_color=COLOR_SKY_100, border_color=COLOR_SKY_300, font_size=10, font_color=COLOR_TEXT_TITLE, bold=True)
    draw_card(slide, "店", Inches(5.95), Inches(4.5), Inches(1.0), Inches(0.55), 
              fill_color=COLOR_SKY_100, border_color=COLOR_SKY_300, font_size=10, font_color=COLOR_TEXT_TITLE, bold=True)
    draw_card(slide, "業務", Inches(7.1), Inches(4.5), Inches(1.0), Inches(0.55), 
              fill_color=COLOR_SKY_100, border_color=COLOR_SKY_300, font_size=10, font_color=COLOR_TEXT_TITLE, bold=True)
              
    # C1. 網路底下的 Level 4 (自媒體、官網媒體)
    # 從「網路」往下垂線分叉
    draw_line(slide, Inches(4.45), Inches(5.15), Inches(5.45), Inches(5.15), COLOR_SKY_500, thickness=1.5)
    draw_line(slide, Inches(5.2), Inches(5.05), Inches(5.2), Inches(5.15), COLOR_SKY_500, thickness=1.5)
    
    # 自媒體
    draw_card(slide, "自媒體", Inches(3.9), Inches(5.35), Inches(1.0), Inches(0.55), 
              fill_color=COLOR_SKY_50, border_color=COLOR_SKY_200, font_size=9, font_color=COLOR_SLATE_700, bold=False)
    draw_line(slide, Inches(4.45), Inches(5.15), Inches(4.45), Inches(5.35), COLOR_SKY_500, thickness=1.5)
    
    # 官網 / 媒體
    draw_card(slide, "官網 / 媒體", Inches(5.05), Inches(5.35), Inches(1.2), Inches(0.55), 
              fill_color=COLOR_SKY_50, border_color=COLOR_SKY_200, font_size=9, font_color=COLOR_SLATE_700, bold=False)
    draw_line(slide, Inches(5.45), Inches(5.15), Inches(5.45), Inches(5.35), COLOR_SKY_500, thickness=1.5)

# ==================== 主程式 ====================

def main():
    input_path = r"真靈光企劃書所有資訊(線上共編處)\真靈光企劃書_給 AI 看的.pptx"
    output_path = r"真靈光企劃書所有資訊(線上共編處)\真靈光企劃書_給 AI 看的_排版美編完成.pptx"
    
    if not os.path.exists(input_path):
        print(f"Error: Input file not found: {input_path}")
        sys.exit(1)
        
    print(f"Opening target PPTX for recreation: {input_path}")
    prs = Presentation(input_path)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Processing Slide {slide_idx+1}...")
        
        # 1. 套用微灰白背景 (#FCFCFC)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG
        
        # 2. 刪除 Google Slide 殘留細線外框
        shapes_to_delete = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE.RECTANGLE:
                left = shape.left.inches if shape.left else 0
                top = shape.top.inches if shape.top else 0
                w = shape.width.inches if shape.width else 0
                h = shape.height.inches if shape.height else 0
                if left < 0.6 and top < 0.6 and w > 12.0 and h > 6.5:
                    if not shape.has_text_frame or not any(p.text.strip() for p in shape.text_frame.paragraphs):
                        shapes_to_delete.append(shape)
        for shape in shapes_to_delete:
            try:
                shape.element.getparent().remove(shape.element)
            except Exception:
                pass
                
        # 3. 尋找與重繪標題
        title_shape = None
        if slide.shapes.title:
            title_shape = slide.shapes.title
        else:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    top = shape.top.inches if shape.top else 99
                    text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
                    if top < 1.4 and text:
                        title_shape = shape
                        break
                        
        if title_shape and title_shape.has_text_frame:
            title_shape.left = Inches(1.05)
            title_shape.top = Inches(0.60)
            title_shape.width = Inches(11.48)
            title_shape.height = Inches(0.8)
            
            tf = title_shape.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "微軟正黑體"
                    run.font.size = Pt(32)
                    run.font.bold = True
                    run.font.color.rgb = COLOR_TEXT_TITLE
                    
            # 畫天空藍裝飾條 (Accent Bar)
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.65), Inches(0.08), Inches(0.6))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = COLOR_ACCENT_BAR
            accent_bar.line.fill.background()
            
            # 畫底部分割線 (Divider Line)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.01))
            line.fill.solid()
            line.fill.fore_color.rgb = COLOR_DIVIDER
            line.line.fill.background()
            
        # 4. 對內容元件套用美化或全新繪製
        # 分流處理特定 Slide：
        if slide_idx == 4:
            # Slide 5: 重新繪製身心靈流程
            print("  Re-drawing Slide 5 (身心靈流程) from Mermaid structure...")
            recreate_slide_5(slide, title_shape)
            
        elif slide_idx == 5:
            # Slide 6: 半自動化循環流程 (維持原有的樣式調整即可，不用重繪)
            print("  Styling Slide 6...")
            for shape in list(slide.shapes):
                process_slide_6_shapes(shape, title_shape)
                
        elif slide_idx == 6:
            # Slide 7: 重新繪製一站式培訓系統
            print("  Re-drawing Slide 7 (一站式培訓系統) from Mermaid structure...")
            recreate_slide_7(slide, title_shape)
            
        elif slide_idx == 9:
            # Slide 10: 重新繪製五量理論
            print("  Re-drawing Slide 10 (五量理論) from Mermaid structure...")
            recreate_slide_10(slide, title_shape)
            
        elif slide_idx == 10:
            # Slide 11: 重新繪製企業架構
            print("  Re-drawing Slide 11 (企業架構) from Mermaid structure...")
            recreate_slide_11(slide, title_shape)
            
        else:
            # 其他一般或表格簡報頁面
            for shape in list(slide.shapes):
                if shape.has_table:
                    table = shape.table
                    for row_i, row in enumerate(table.rows):
                        for cell in row.cells:
                            if row_i == 0:
                                apply_shape_styling(cell, fill_rgb=COLOR_SKY_600)
                                if cell.text_frame:
                                    format_text_frame(cell.text_frame, default_size=13, default_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
                            else:
                                apply_shape_styling(cell, fill_rgb=COLOR_SKY_50)
                                if cell.text_frame:
                                    format_text_frame(cell.text_frame, default_size=11, default_color=COLOR_SLATE_700, bold=False, align=PP_ALIGN.CENTER)
                else:
                    # 一般 shapes 遞迴格式化字型
                    # 為防止有子 shapes 因寫死文字重置而遺失，此處不破壞非重繪頁面的 shapes，僅修正其 font/color
                    process_general_shape(shape, title_shape)
                    
    # 5. 儲存
    prs.save(output_path)
    print(f"\nRebuilt and recreated presentation saved successfully to: {output_path}")

# ==================== 原有的樣式修正輔助函數 ====================

def process_slide_6_shapes(shape, title_shape):
    if shape == title_shape:
        return
    shape_type = shape.shape_type
    if shape_type == 6 or hasattr(shape, "shapes"):
        try:
            for subshape in shape.shapes:
                process_slide_6_shapes(subshape, title_shape)
        except Exception:
            pass
        return
    text = ""
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
    if text in ["宣傳", "量測", "課程", "產品"]:
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_100, line_rgb=COLOR_SKY_300, line_width=1.5)
        format_text_frame(shape.text_frame, default_size=18, default_color=COLOR_TEXT_TITLE, bold=True, align=PP_ALIGN.CENTER)
        return
    if (shape_type == MSO_SHAPE.RECTANGLE or shape_type == MSO_SHAPE.ROUNDED_RECTANGLE) and not text:
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_50, line_rgb=COLOR_SKY_200, line_width=1.5)
        return
    is_arrow = (shape_type in [MSO_SHAPE.RIGHT_ARROW, MSO_SHAPE.LEFT_ARROW, MSO_SHAPE.UP_ARROW, MSO_SHAPE.DOWN_ARROW])
    is_small_freeform = (shape_type == 5 and not text and shape.width and shape.width.inches < 1.0)
    if is_arrow or is_small_freeform:
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_500, line_rgb=COLOR_SKY_500)
        return
    if shape.has_text_frame and text:
        align = PP_ALIGN.LEFT
        if len(text) < 5:
            align = PP_ALIGN.CENTER
        size = 12
        if len(text) > 20:
            size = 11
        format_text_frame(shape.text_frame, default_size=size, default_color=COLOR_SLATE_700, bold=False, align=align)

def process_general_shape(shape, title_shape):
    if shape == title_shape:
        return
    shape_type = shape.shape_type
    if shape_type == 6 or hasattr(shape, "shapes"):
        try:
            for subshape in shape.shapes:
                process_general_shape(subshape, title_shape)
        except Exception:
            pass
        return
    text = ""
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
    if text:
        align = PP_ALIGN.LEFT
        if len(text) < 8:
            align = PP_ALIGN.CENTER
        format_text_frame(shape.text_frame, default_size=16, default_color=COLOR_SLATE_700, bold=False, align=align)
    is_arrow = (shape_type in [MSO_SHAPE.RIGHT_ARROW, MSO_SHAPE.LEFT_ARROW, MSO_SHAPE.UP_ARROW, MSO_SHAPE.DOWN_ARROW])
    if is_arrow:
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_500, line_rgb=COLOR_SKY_500)

if __name__ == "__main__":
    main()
