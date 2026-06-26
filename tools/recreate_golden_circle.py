import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml import parse_xml

def add_arrow_line(slide, x1, y1, x2, y2, color, width_pt=2):
    # Add straight connector
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line = connector.line
    line.color.rgb = color
    line.width = Pt(width_pt)
    
    # Append XML to add arrowhead at the end (tailEnd)
    line_elem = line._get_or_add_ln()
    tail_end_xml = '<a:tailEnd type="triangle" w="med" len="med" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    line_elem.append(parse_xml(tail_end_xml))
    return connector

def recreate_golden_circle():
    pptx_path = r"d:\Git\true-soul-light-cowork\最終輸出用檔案\重製.pptx"
    
    if not os.path.exists(pptx_path):
        print(f"Error: {pptx_path} does not exist.")
        sys.exit(1)
        
    prs = Presentation(pptx_path)
    
    # Ensure slide size is 13.33 x 7.50 inches (Widescreen 16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.50)
    
    slide = prs.slides[0]
    
    # Delete all existing shapes (including the original image screenshot)
    shapes_to_delete = list(slide.shapes)
    for shape in shapes_to_delete:
        el = shape.element
        el.getparent().remove(el)
        
    # 1. Apply slide theme background color (#FCFCFC)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(252, 252, 252)
    
    # 2. Add Title decoration vertical bar (Sky-500: #0EA5E9)
    # Position: Left=0.80", Top=0.65", Width=0.08", Height=0.60"
    rect_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.80), Inches(0.65), Inches(0.08), Inches(0.60)
    )
    rect_bar.fill.solid()
    rect_bar.fill.fore_color.rgb = RGBColor(14, 165, 233)
    rect_bar.line.fill.background() # No border
    
    # 3. Add Title text box (Microsoft JhengHei, 32pt, Bold, Deep Cyan: #083344)
    # Position: Left=1.05", Top=0.60", Width=11.48", Height=0.80"
    title_box = slide.shapes.add_textbox(
        Inches(1.05), Inches(0.60), Inches(11.48), Inches(0.80)
    )
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = Inches(0)
    tf_title.margin_right = Inches(0)
    tf_title.margin_top = Inches(0)
    tf_title.margin_bottom = Inches(0)
    
    p_title = tf_title.paragraphs[0]
    p_title.text = "行銷的底層邏輯：黃金圈法則 (Golden Circle)"
    p_title.font.name = "微軟正黑體"
    p_title.font.size = Pt(32)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(8, 51, 68)
    
    # 4. Add horizontal divider line (Slate-200: #E2E8F0)
    # Position: Left=0.80", Top=1.45", Width=11.73", Height=0.01"
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.80), Inches(1.45), Inches(11.73), Inches(0.01)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(226, 232, 240)
    divider.line.fill.background()
    
    # 5. Draw concentric circles (What -> How -> Why)
    # Center = (3.50", 4.40")
    
    # 最外圈 (What)
    # Radius = 2.20" -> Diameter = 4.40"
    # Left = 3.50 - 2.20 = 1.30", Top = 4.40 - 2.20 = 2.20"
    what_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.30), Inches(2.20), Inches(4.40), Inches(4.40)
    )
    what_circle.fill.solid()
    what_circle.fill.fore_color.rgb = RGBColor(254, 240, 138) # #FEF08A (Yellow-200)
    what_circle.line.fill.background()
    
    # 中間圈 (How)
    # Radius = 1.50" -> Diameter = 3.00"
    # Left = 3.50 - 1.50 = 2.00", Top = 4.40 - 1.50 = 2.90"
    how_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(2.00), Inches(2.90), Inches(3.00), Inches(3.00)
    )
    how_circle.fill.solid()
    how_circle.fill.fore_color.rgb = RGBColor(253, 186, 116) # #FDBA74 (Orange-300)
    how_circle.line.fill.background()
    
    # 最內圈 (Why)
    # Radius = 0.80" -> Diameter = 1.60"
    # Left = 3.50 - 0.80 = 2.70", Top = 4.40 - 0.80 = 3.60"
    why_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(2.70), Inches(3.60), Inches(1.60), Inches(1.60)
    )
    why_circle.fill.solid()
    why_circle.fill.fore_color.rgb = RGBColor(249, 115, 22) # #F97316 (Orange-500)
    why_circle.line.fill.background()
    
    # Add text "Why" inside the innermost circle
    tf_why = why_circle.text_frame
    tf_why.word_wrap = True
    p_why = tf_why.paragraphs[0]
    p_why.alignment = PP_ALIGN.CENTER
    p_why.text = "Why"
    p_why.font.name = "微軟正黑體"
    p_why.font.size = Pt(24)
    p_why.font.bold = True
    p_why.font.color.rgb = RGBColor(255, 255, 255)
    
    # 6. Add text "How" and "What" on the circles
    # How text (in the How ring, below Why circle)
    # Position: Left=2.00", Top=5.20", Width=3.00", Height=0.50"
    how_text_box = slide.shapes.add_textbox(
        Inches(2.00), Inches(5.20), Inches(3.00), Inches(0.50)
    )
    tf_how = how_text_box.text_frame
    tf_how.word_wrap = True
    p_how = tf_how.paragraphs[0]
    p_how.alignment = PP_ALIGN.CENTER
    p_how.text = "How"
    p_how.font.name = "微軟正黑體"
    p_how.font.size = Pt(20)
    p_how.font.bold = True
    p_how.font.color.rgb = RGBColor(255, 255, 255)
    
    # What text (in the What ring, below How circle)
    # Position: Left=1.30", Top=5.95", Width=4.40", Height=0.50"
    what_text_box = slide.shapes.add_textbox(
        Inches(1.30), Inches(5.95), Inches(4.40), Inches(0.50)
    )
    tf_what = what_text_box.text_frame
    tf_what.word_wrap = True
    p_what = tf_what.paragraphs[0]
    p_what.alignment = PP_ALIGN.CENTER
    p_what.text = "What"
    p_what.font.name = "微軟正黑體"
    p_what.font.size = Pt(20)
    p_what.font.bold = True
    p_what.font.color.rgb = RGBColor(255, 255, 255)
    
    # 7. Add Right descriptions
    # Why description: Left=8.50", Top=2.10", Width=4.00", Height=1.20"
    why_desc = slide.shapes.add_textbox(
        Inches(8.50), Inches(2.10), Inches(4.00), Inches(1.20)
    )
    tf_why_desc = why_desc.text_frame
    tf_why_desc.word_wrap = True
    tf_why_desc.margin_left = Inches(0)
    
    p_why_title = tf_why_desc.paragraphs[0]
    p_why_title.text = "為何而做"
    p_why_title.font.name = "微軟正黑體"
    p_why_title.font.size = Pt(28)
    p_why_title.font.bold = True
    p_why_title.font.color.rgb = RGBColor(234, 88, 12) # #EA580C (Orange-600)
    
    p_why_detail = tf_why_desc.add_paragraph()
    p_why_detail.text = "動機、初衷與目的"
    p_why_detail.font.name = "微軟正黑體"
    p_why_detail.font.size = Pt(20)
    p_why_detail.font.color.rgb = RGBColor(30, 41, 59) # #1E293B
    p_why_detail.space_before = Pt(8)
    
    # How description: Left=8.50", Top=3.80", Width=4.00", Height=1.20"
    how_desc = slide.shapes.add_textbox(
        Inches(8.50), Inches(3.80), Inches(4.00), Inches(1.20)
    )
    tf_how_desc = how_desc.text_frame
    tf_how_desc.word_wrap = True
    tf_how_desc.margin_left = Inches(0)
    
    p_how_title = tf_how_desc.paragraphs[0]
    p_how_title.text = "如何執行"
    p_how_title.font.name = "微軟正黑體"
    p_how_title.font.size = Pt(28)
    p_how_title.font.bold = True
    p_how_title.font.color.rgb = RGBColor(217, 119, 6) # #D97706 (Amber-600)
    
    p_how_detail = tf_how_desc.add_paragraph()
    p_how_detail.text = "策略、方法與流程"
    p_how_detail.font.name = "微軟正黑體"
    p_how_detail.font.size = Pt(20)
    p_how_detail.font.color.rgb = RGBColor(30, 41, 59)
    p_how_detail.space_before = Pt(8)
    
    # What description: Left=8.50", Top=5.50", Width=4.00", Height=1.20"
    what_desc = slide.shapes.add_textbox(
        Inches(8.50), Inches(5.50), Inches(4.00), Inches(1.20)
    )
    tf_what_desc = what_desc.text_frame
    tf_what_desc.word_wrap = True
    tf_what_desc.margin_left = Inches(0)
    
    p_what_title = tf_what_desc.paragraphs[0]
    p_what_title.text = "成果為何"
    p_what_title.font.name = "微軟正黑體"
    p_what_title.font.size = Pt(28)
    p_what_title.font.bold = True
    p_what_title.font.color.rgb = RGBColor(161, 98, 7) # #A16207 (Yellow-700)
    
    p_what_detail = tf_what_desc.add_paragraph()
    p_what_detail.text = "成果面的體現"
    p_what_detail.font.name = "微軟正黑體"
    p_what_detail.font.size = Pt(20)
    p_what_detail.font.color.rgb = RGBColor(30, 41, 59)
    p_what_detail.space_before = Pt(8)
    
    # 8. Draw arrow lines with XML arrowhead
    # Why arrow: Start=(4.06, 3.84), End=(8.30, 2.50)
    add_arrow_line(slide, 4.06, 3.84, 8.30, 2.50, RGBColor(249, 115, 22), width_pt=2)
    # How arrow: Start=(5.00, 4.40), End=(8.30, 4.40)
    add_arrow_line(slide, 5.00, 4.40, 8.30, 4.40, RGBColor(253, 186, 116), width_pt=2)
    # What arrow: Start=(5.06, 5.96), End=(8.30, 6.00)
    add_arrow_line(slide, 5.06, 5.96, 8.30, 6.00, RGBColor(254, 240, 138), width_pt=2)
    
    # Save presentation
    try:
        prs.save(pptx_path)
        print("Successfully saved directly to 重製.pptx")
    except PermissionError:
        # Fallback to _v2 if PPTX is open/locked by the user
        dirname = os.path.dirname(pptx_path)
        fallback_path = os.path.join(dirname, "重製_排版美編完成.pptx")
        prs.save(fallback_path)
        print(f"PermissionError: file locked. Saved to fallback path instead: {fallback_path}")

if __name__ == "__main__":
    recreate_golden_circle()
