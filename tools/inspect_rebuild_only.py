import os
import sys
from pptx import Presentation

# 避免中文編碼問題，輸出用 utf-8
sys.stdout.reconfigure(encoding='utf-8')

def inspect_pptx_full(path, output_log):
    with open(output_log, 'w', encoding='utf-8') as f:
        f.write(f"Inspecting PPTX: {path}\n")
        if not os.path.exists(path):
            f.write("File not found!\n")
            return
        
        try:
            prs = Presentation(path)
            f.write(f"Total slides: {len(prs.slides)}\n")
            f.write(f"Slide Size: Width={prs.slide_width.inches:.3f}\", Height={prs.slide_height.inches:.3f}\"\n")
            
            for i, slide in enumerate(prs.slides):
                f.write(f"\n--- Slide {i+1} ---\n")
                
                # Check background
                bg = slide.background
                if bg and bg.fill:
                    f.write(f"  Background Fill Type: {bg.fill.type}\n")
                    if bg.fill.type == 1: # Solid
                        f.write(f"  Background Color: {bg.fill.fore_color.rgb if bg.fill.fore_color else 'None'}\n")
                
                # Check shapes
                for shape in slide.shapes:
                    shape_type = shape.shape_type
                    shape_name = shape.name
                    left = shape.left.inches if shape.left else 0
                    top = shape.top.inches if shape.top else 0
                    width = shape.width.inches if shape.width else 0
                    height = shape.height.inches if shape.height else 0
                    
                    text_preview = ""
                    if shape.has_text_frame:
                        texts = []
                        for p in shape.text_frame.paragraphs:
                            if p.text.strip():
                                texts.append(p.text.strip())
                        text_preview = " | Texts: " + " / ".join(texts)
                    
                    f.write(f"  Shape: '{shape_name}', Type: {shape_type}, Pos: L={left:.2f}\", T={top:.2f}\", W={width:.2f}\", H={height:.2f}\"{text_preview}\n")
            print(f"Inspection complete. Log saved to {output_log}")
        except Exception as e:
            f.write(f"Error during inspection: {str(e)}\n")
            print(f"Error during inspection: {str(e)}")

# Inspect Target PPTX
inspect_pptx_full(r"真靈光企劃書所有資訊(線上共編處)\重製.pptx", "tools/rebuild_target_inspection.txt")

# Inspect Style Reference
inspect_pptx_full(r"最終輸出用檔案\舊版、已使用封存\行銷界九陽神功_完整講稿_總監編修#4.pptx", "tools/style_reference_inspection.txt")
