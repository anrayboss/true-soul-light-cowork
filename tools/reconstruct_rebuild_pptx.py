import os
import sys
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 避免中文編碼輸出問題
sys.stdout.reconfigure(encoding='utf-8')

# 定義配色系統的 RGBColor 物件
COLOR_BG = RGBColor(252, 252, 252)        # #FCFCFC (九陽神功背景色)
COLOR_TEXT_TITLE = RGBColor(8, 51, 68)     # #083344 (Cyan-950 標題文字)
COLOR_ACCENT_BAR = RGBColor(14, 165, 233)  # #0EA5E9 (Sky-500 裝飾條)
COLOR_DIVIDER = RGBColor(226, 232, 240)    # #E2E8F0 (Slate-200 分割線)

COLOR_SKY_50 = RGBColor(240, 249, 255)     # #F0F9FF (主要卡片底色)
COLOR_SKY_100 = RGBColor(224, 242, 254)    # #E0F2FE (主流程/次標底色)
COLOR_SKY_200 = RGBColor(186, 230, 253)    # #BAE6FD (卡片細邊框)
COLOR_SKY_300 = RGBColor(125, 211, 252)    # #7DD3FC (次要流程邊框)
COLOR_SKY_500 = RGBColor(14, 165, 233)     # #0EA5E9 (主流程底色/箭頭)
COLOR_SKY_600 = RGBColor(2, 132, 199)      # #0284C7 (大分類底色)
COLOR_SKY_700 = RGBColor(3, 105, 161)      # #0369A1 (頂部核心底色)

COLOR_SLATE_700 = RGBColor(51, 65, 85)     # #334155 (主文字)
COLOR_SLATE_500 = RGBColor(100, 116, 139)   # #64748B (輔助說明文字)
COLOR_WHITE = RGBColor(255, 255, 255)      # 白色

def format_text_frame(tf, font_name="微軟正黑體", default_size=12, default_color=COLOR_SLATE_700, bold=False, align=None):
    """
    統一格式化文字框內部的所有段落與 Run，防止出現新細明體，並調整對齊與文字大小。
    """
    tf.word_wrap = True
    for p in tf.paragraphs:
        if align is not None:
            p.alignment = align
        p.font.name = font_name
        
        # 如果段落有文字
        if p.text.strip():
            # 遍歷 Runs
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(default_size)
                run.font.color.rgb = default_color
                run.font.bold = bold

def apply_shape_styling(shape, fill_rgb=None, line_rgb=None, line_width=1.5):
    """
    安全地套用填充色與邊框顏色到 shape 物件。
    """
    try:
        if fill_rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
        else:
            # 透明背景
            shape.fill.background()
            
        if hasattr(shape, 'line') and shape.line:
            if line_rgb:
                shape.line.color.rgb = line_rgb
                shape.line.width = Pt(line_width)
            else:
                # 若無指定邊框，將邊框設成與填充相同，以隱藏邊框
                if fill_rgb:
                    shape.line.color.rgb = fill_rgb
                else:
                    shape.line.fill.background()
    except Exception as e:
        pass

def process_slide_1_shapes(shape, title_shape):
    """
    專屬 Slide 1 (半自動化循環流程) 的形狀美化邏輯
    """
    if shape == title_shape:
        return
        
    shape_type = shape.shape_type
    shape_name = shape.name
    
    # 遞迴處理 Group
    if shape_type == 6 or hasattr(shape, "shapes"):
        try:
            for subshape in shape.shapes:
                process_slide_1_shapes(subshape, title_shape)
        except Exception as e:
            print(f"Error recursing Group on Slide 1: {e}")
        return
        
    text = ""
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
        
    # 1. 核心四個流程圓形 ("宣傳"、"量測"、"課程"、"產品")
    if text in ["宣傳", "量測", "課程", "產品"]:
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_100, line_rgb=COLOR_SKY_300, line_width=1.5)
        format_text_frame(shape.text_frame, default_size=18, default_color=COLOR_TEXT_TITLE, bold=True, align=PP_ALIGN.CENTER)
        return
        
    # 2. 說明文字大背景卡片 (圓角矩形/矩形)
    # 我們可以根據 shape 寬度與高度特徵，或者若它不含文字但屬於卡片框：
    if (shape_type == MSO_SHAPE.RECTANGLE or shape_type == MSO_SHAPE.ROUNDED_RECTANGLE) and not text:
        # Slide 1 中的 4 個大卡片背景：331, 329, 333, 335
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_50, line_rgb=COLOR_SKY_200, line_width=1.5)
        return

    # 3. 箭頭形狀或是沒有文字的流程箭頭 (如 345, 346 以及其他連接線形狀)
    is_arrow = (shape_type in [MSO_SHAPE.RIGHT_ARROW, MSO_SHAPE.LEFT_ARROW, MSO_SHAPE.UP_ARROW, MSO_SHAPE.DOWN_ARROW])
    is_small_freeform_flow = (shape_type == 5 and not text and shape.width and shape.width.inches < 1.0)
    if is_arrow or is_small_freeform_flow:
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_500, line_rgb=COLOR_SKY_500)
        return
        
    # 4. 一般說明文字框
    if shape.has_text_frame and text:
        # 例如 "健康食品 / 美容 / 保養品..."
        # 我們將文字對齊左側，字體設為 12pt，顏色為 Slate-700
        # 如果是 "宣傳", "量測" 底下的圓圈文字本身已經在 (1) 處理，這裡處理大卡片文字與其他零散文字
        align = PP_ALIGN.LEFT
        # 如果是流程圖圓圈上面的小字，或是位置偏置中的文字，則置中
        if len(text) < 5:
            align = PP_ALIGN.CENTER
            
        size = 12
        # 若是比較長的條目文字，稍小一點以便美觀
        if len(text) > 20:
            size = 11
            
        format_text_frame(shape.text_frame, default_size=size, default_color=COLOR_SLATE_700, bold=False, align=align)

def process_slide_2_shapes(shape, title_shape):
    """
    專屬 Slide 2 (身心靈一站式服務培訓系統) 的形狀美化邏輯
    """
    if shape == title_shape:
        return
        
    shape_type = shape.shape_type
    shape_name = shape.name
    
    # 遞迴處理 Group
    if shape_type == 6 or hasattr(shape, "shapes"):
        try:
            for subshape in shape.shapes:
                process_slide_2_shapes(subshape, title_shape)
        except Exception as e:
            print(f"Error recursing Group on Slide 2: {e}")
        return
        
    text = ""
    if shape.has_text_frame:
        text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
        
    # 根據 Slide 2 的元件文字內容和形狀類型做精細美化：
    
    # 1. 橫向分類大標籤 ("客戶流程", "培訓系統")
    if text in ["客戶流程", "培訓系統"]:
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_600, line_rgb=None)
        format_text_frame(shape.text_frame, default_size=18, default_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
        return
        
    # 2. 頂部核心 ("身心靈整合平台")
    if text == "身心靈整合平台":
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_700, line_rgb=None)
        format_text_frame(shape.text_frame, default_size=16, default_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
        return
        
    # 3. 客戶流程階段核心 ("宣傳", "量測", "諮詢/課程", "產品")
    if text in ["宣傳", "量測", "諮詢/課程", "產品"]:
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_500, line_rgb=None)
        format_text_frame(shape.text_frame, default_size=16, default_color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
        return
        
    # 4. 培訓系統階段核心 ("行銷與業務培訓", "AI身心靈檢測", "身心靈老師培訓", "產品培訓")
    if text in ["行銷與業務培訓", "AI身心靈檢測", "身心靈老師培訓", "產品培訓"]:
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_100, line_rgb=COLOR_SKY_300, line_width=1.5)
        format_text_frame(shape.text_frame, default_size=14, default_color=COLOR_TEXT_TITLE, bold=True, align=PP_ALIGN.CENTER)
        return
        
    # 5. 說明小卡片 (如 "整合一致性"、"流量"、"流量變現"、"可數據化"、"身心靈量測"、"獨特技術諮詢"、"可複製靈性提升"、"建議顧客所需"、"獨特產品")
    # 這些卡片在 inspection_log 中有 text_box 與對應的 auto_shape 背景。
    # 如果是有文字的卡片：
    if text:
        # 如果是說明文字框，本身不帶背景的，我們將文字格式化為 Slate-700, 12pt
        # 如果它是背景卡片 shape，則下面 (6) 會處理
        format_text_frame(shape.text_frame, default_size=12, default_color=COLOR_SLATE_700, bold=False, align=PP_ALIGN.CENTER)
        return
        
    # 6. 無文字的背景框 (通常是卡片的背景 auto_shape)
    if (shape_type == MSO_SHAPE.RECTANGLE or shape_type == MSO_SHAPE.ROUNDED_RECTANGLE) and not text:
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_50, line_rgb=COLOR_SKY_200, line_width=1.5)
        return
        
    # 7. 流程箭頭或小連線形狀 (無文字且寬高較扁或 shape_type 為箭頭或線)
    is_arrow = (shape_type in [MSO_SHAPE.RIGHT_ARROW, MSO_SHAPE.LEFT_ARROW, MSO_SHAPE.UP_ARROW, MSO_SHAPE.DOWN_ARROW, MSO_SHAPE.UP_DOWN_ARROW])
    is_line = (shape_type == 9 or shape_name.startswith("直線") or shape_name.startswith("Connector"))
    is_small_freeform = (shape_type == 5 and shape.width and shape.width.inches < 1.0)
    
    if is_arrow or is_line or is_small_freeform:
        apply_shape_styling(shape, fill_rgb=COLOR_SKY_500, line_rgb=COLOR_SKY_500)
        return

def main():
    input_path = r"真靈光企劃書所有資訊(線上共編處)\重製.pptx"
    output_path = r"真靈光企劃書所有資訊(線上共編處)\重製_排版美編完成.pptx"
    
    if not os.path.exists(input_path):
        print(f"Error: Input file not found: {input_path}")
        sys.exit(1)
        
    print(f"Opening target PPTX: {input_path}")
    prs = Presentation(input_path)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"Processing Slide {slide_idx+1}...")
        
        # 1. 套用高級微灰白背景 (#FCFCFC)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG
        
        # 2. 刪除任何 Google Slide 匯出時殘留的外層邊框
        shapes_to_delete = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE.RECTANGLE:
                left = shape.left.inches if shape.left else 0
                top = shape.top.inches if shape.top else 0
                w = shape.width.inches if shape.width else 0
                h = shape.height.inches if shape.height else 0
                # 外框特徵：寬高涵蓋幾乎整張投影片，且無文字
                if left < 0.6 and top < 0.6 and w > 12.0 and h > 6.5:
                    if not shape.has_text_frame or not any(p.text.strip() for p in shape.text_frame.paragraphs):
                        shapes_to_delete.append(shape)
                        
        for shape in shapes_to_delete:
            try:
                shape.element.getparent().remove(shape.element)
                print(f"  Removed outer frame border.")
            except Exception as e:
                print(f"  Failed to remove outer border: {e}")
                
        # 3. 尋找與重塑投影片標題
        title_shape = None
        if slide.shapes.title:
            title_shape = slide.shapes.title
        else:
            # 尋找頂部位置像標題的文字框
            for shape in slide.shapes:
                if shape.has_text_frame:
                    top = shape.top.inches if shape.top else 99
                    text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
                    if top < 1.4 and text:
                        title_shape = shape
                        break
                        
        if title_shape and title_shape.has_text_frame:
            # 移動到九陽神功標準標題位置
            title_shape.left = Inches(1.05)
            title_shape.top = Inches(0.60)
            title_shape.width = Inches(11.48)
            title_shape.height = Inches(0.8)
            
            tf = title_shape.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            
            # 套用標題文字樣式
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "微軟正黑體"
                    run.font.size = Pt(32)
                    run.font.bold = True
                    run.font.color.rgb = COLOR_TEXT_TITLE
                    
            # 畫天空藍裝飾條 (Accent Bar)
            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8), Inches(0.65), Inches(0.08), Inches(0.6)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = COLOR_ACCENT_BAR
            accent_bar.line.fill.background()
            
            # 畫底部分割線 (Divider Line)
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.01)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = COLOR_DIVIDER
            line.line.fill.background()
            
            print(f"  Styled slide title and added accent/divider visual elements.")
            
        # 4. 對內容元件套用九陽風格的色彩與排版
        if slide_idx == 0:
            # 處理投影片一
            for shape in list(slide.shapes):
                process_slide_1_shapes(shape, title_shape)
        elif slide_idx == 1:
            # 處理投影片二
            for shape in list(slide.shapes):
                process_slide_2_shapes(shape, title_shape)
                
    # 5. 儲存重製簡報
    prs.save(output_path)
    print(f"\nRebuilt presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    main()
